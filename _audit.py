import subprocess, sys, glob, os, re
from pptx import Presentation
from pptx.util import Emu
from src.services.verify_sources import verify_ticker
TICKERS = ["1180.SR","BKMB.OM","2010.SR","2020.SR","ORDS.QA","9988.HK","0981.HK"]
def latest(tk):
    fs=sorted(glob.glob(f"outputs/{tk}_2026*.pptx")); return fs[-1] if fs else None
def cells(slide, ylo, yhi):
    return [(round(Emu(sh.top).inches,2),Emu(sh.left).inches,sh.text_frame.text.strip())
            for sh in slide.shapes if sh.has_text_frame and sh.text_frame.text.strip()
            and ylo<Emu(sh.top).inches<yhi]
def audit(tk):
    f=latest(tk)
    if not f: return "NO DECK"
    p=Presentation(f); s=list(p.slides)
    out={}
    # Slide 1 key data (y~4.98) + consensus (y~3.64) + perf (y~6.05) + 52w
    kd=[t for y,x,t in sorted(cells(s[0],4.9,5.1))]
    cons=[t for y,x,t in sorted(cells(s[0],3.55,3.75))]
    perf=[t for y,x,t in sorted(cells(s[0],5.95,6.15))]
    rng=[t for y,x,t in sorted(cells(s[0],6.8,7.0))]
    out['S1 consensus(rating/tgt/upside)']=cons
    out['S1 keydata(close/mcap/next/PE/div/ccy)']=kd
    out['S1 perf(1D..YTD)']=perf
    out['S1 52w']=rng
    # Slide 2 expectations table rows
    rows={}
    for y,x,t in cells(s[1],4.7,6.0): rows.setdefault(y,[]).append((x,t))
    tbl=[' | '.join(t for x,t in sorted(v)) for y,v in sorted(rows.items())]
    out['S2 expectations']=tbl
    # Slide 3 peer table + charts
    prows={}
    for y,x,t in cells(s[2],5.0,8.6): prows.setdefault(y,[]).append((x,t))
    peers=[' | '.join(t for x,t in sorted(v)) for y,v in sorted(prows.items()) if len(v)>=4]
    out['S3 peers']=peers
    out['charts']=sum(1 for sl in s for sh in sl.shapes if sh.shape_type==13)
    out['placeholders']=[sh.text_frame.text[:24] for sl in s for sh in sl.shapes if sh.has_text_frame and sh.text_frame.text.strip().startswith("No ")]
    # count dashes
    blob="\n".join(sh.text_frame.text for sl in s for sh in sl.shapes if sh.has_text_frame)
    out['dash_count']=blob.count("—")
    return out
for tk in TICKERS:
    subprocess.run([sys.executable,"-m","src.main","--ticker",tk,"--mode","preview"],
                   env={**os.environ,"REFRESH_ON_RENDER":"1","DISABLE_LLM":"1"}, capture_output=True, text=True)
    a=audit(tk)
    print(f"\n########## {tk} ##########")
    if isinstance(a,dict):
        for k,v in a.items(): print(f"  {k}: {v}")
        try:
            vr=verify_ticker(tk); vs=vr['summary']
            fails=[(c['field'],c['deck'],c['note'][:50]) for c in vr['checks'] if c['verdict']=='FAIL']
            print(f"  VERIFY: {vs['pass']}p/{vs['warn']}w/{vs['fail']}f  FAILS={fails}")
        except Exception as e: print("  VERIFY ERR",e)
    else: print(" ",a)
print("\nDONE")
