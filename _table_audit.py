import subprocess, sys, glob, os
from pptx import Presentation
from pptx.util import Emu
# Diverse sample across regions/sectors/coverage
SET = ["600519.SS","601398.SS","002594.SZ","300750.SZ","0700.HK","0941.HK","1810.HK",
       "RELIANCE.NS","TCS.NS","HDFCBANK.NS","SUNPHARMA.NS","MARUTI.NS",
       "2222.SR","1120.SR","7010.SR","2010.SR","QNBK.QA","ORDS.QA",
       "NPN.JO","SOL.JO","2222.SR"]
SET=list(dict.fromkeys(SET))
def latest(tk):
    fs=sorted(glob.glob(f"outputs/{tk}_2026*.pptx")); return fs[-1] if fs else None
def audit(tk):
    f=latest(tk)
    if not f: return "NO DECK"
    p=Presentation(f); s=list(p.slides)
    # find table region rows on slide 2
    rows={}
    for sh in s[1].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip():
            y=round(Emu(sh.top).inches,2); x=Emu(sh.left).inches
            if 4.4<y<6.2: rows.setdefault(y,[]).append((x,sh.text_frame.text.strip()))
    header=None; metric_rows=[]
    for y in sorted(rows):
        c=[t for _,t in sorted(rows[y])]
        if c and c[0]=="METRIC": header=c
        elif c and "(" in c[0] and len(c)>=4: metric_rows.append(c)
    if not header: return "no table"
    mode = "ANNUAL" if "CAGR" in " ".join(header) else "QUARTERLY"
    # count rows with >=1 real numeric value (not all dashes)
    def has_val(c): return any(x not in ("—","n/m","") and any(ch.isdigit() for ch in x) for x in c[1:])
    filled=sum(1 for c in metric_rows if has_val(c))
    return f"{mode} rows={len(metric_rows)} filled={filled}"
for tk in SET:
    subprocess.run([sys.executable,"-m","src.main","--ticker",tk,"--mode","preview"],
                   env={**os.environ,"REFRESH_ON_RENDER":"1","DISABLE_LLM":"1"}, capture_output=True, text=True)
    print(f"{tk:12} {audit(tk)}", flush=True)
print("DONE")
