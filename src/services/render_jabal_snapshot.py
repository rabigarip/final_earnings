"""
Jabal Asset Management — Slide 1 (Snapshot) renderer.

Renders the first slide of the new 3-slide preview deck following the
spec in `docs/stage2/design_spec.md`. Reads exclusively from
`canonical_store` — never calls a provider directly. Returns the
pptx Slide object so a higher-level builder can stack slides 1/2/3.

Layout (top→bottom):
  1. Header strip
  2. Title block (kicker, company name, meta line, period subtitle)
  3. Analyst consensus row (3 cards)
  4. Key data row (6 metric blocks)
  5. Recent performance row (6 colored deltas)
  6. 52-week range bar
  7. Analyst highlights (5 pill rows)
  8. Footer
"""

from __future__ import annotations

from dataclasses import dataclass
from datetime import datetime
from typing import Optional

import math
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.util import Inches, Pt, Emu

from src.services.jabal_design_tokens import (
    BLACK, GRAY, MUTED, GOLD, POS, NEG, CARD, WHITE,
    FONT_DISPLAY, FONT_UI,
    SZ_HERO, SZ_KICKER, SZ_VALUE, SZ_VALUE_LG, SZ_LABEL, SZ_BODY,
    SZ_META, SZ_HEADER, SZ_FOOTER, SZ_BULLET_PILL, SZ_TAB_NUM,
    PAGE_W_IN, PAGE_H_IN, MARGIN_L, MARGIN_R, CONTENT_W,
    RULE_THICK_PT, BORDER_THICK_PT, LEFT_ACCENT_W_IN,
    in_, signed_color,
)
from src.services.canonical_store import get_all_fields, CanonicalValue


def _normalize_date(s: str) -> str:
    """Coerce ISO / US-short / month-name strings to '%d %b %Y' ('15 Jul 2026')
    so REPORT DATE, broker-action dates and the footer match. Returns the
    input unchanged when no format is recognised."""
    if not s or not isinstance(s, str):
        return s or ""
    s = s.strip()
    if s.upper() == "TBA":
        return s
    from datetime import datetime as _dt
    for fmt in ("%Y-%m-%d", "%m/%d/%y", "%m/%d/%Y", "%d %b %Y", "%d %B %Y",
                "%b. %d, %Y", "%b %d, %Y", "%B %d, %Y", "%b %d", "%b. %d"):
        try:
            d = _dt.strptime(s, fmt)
            if d.year == 1900:  # month-day only → assume current year
                d = d.replace(year=_dt.now().year)
            return d.strftime("%d %b %Y")
        except ValueError:
            continue
    return s


# ── Low-level primitives ────────────────────────────────────

def _fit_pt(text: str, width_in: float, height_in: float, base_pt: float,
            *, wrap: bool, all_caps: bool) -> float:
    """Return a font size (pt) that makes `text` fit `width_in × height_in`,
    shrinking from `base_pt` only when it would overflow. Baked into the
    saved file so it renders correctly in EVERY viewer, not just PowerPoint.

    - A space-less token (a number like "$106.7B", a single word) is never
      wrapped — it's fitted to the width on one line so it can't break and
      misalign a column.
    - A short box (≈ one line tall) is treated as single-line.
    - Multi-word text in a tall box is fitted by wrapped-line height.
    Conservative: a small tolerance and a floor (~60% of base) keep it from
    over-shrinking text that already fits."""
    t = (text or "").strip()
    if not t or width_in <= 0 or height_in <= 0:
        return base_pt
    # caps and digits are wider than mixed-case lowercase
    cw_factor = 0.58 if all_caps else 0.50
    base_lh = base_pt * 1.18 / 72.0
    single_line = (not wrap) or (" " not in t) or (height_in < 1.7 * base_lh)
    size = float(base_pt)
    # Allow a touch more shrink for long LLM prose (catalysts/risks/watch),
    # which vary in length and otherwise clip the bottom of their box. 5.5pt
    # is still legible at the deck's print scale; below that we'd rather the
    # text shrink than overrun the next section.
    floor = max(5.5, base_pt * 0.55)
    while size > floor:
        cw = size * cw_factor / 72.0
        lh = size * 1.18 / 72.0
        if single_line:
            if len(t) * cw <= width_in * 1.03:
                break
        else:
            cpl = max(1, int(width_in / cw))
            lines = max(1, math.ceil(len(t) / cpl))
            if lines * lh <= height_in * 1.03:
                break
        size -= 0.5
    return round(size, 1)


def _text(slide, left, top, width, height, text, *,
          font=FONT_UI, size=SZ_BODY, bold=False, italic=False, color=BLACK,
          align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, all_caps=False,
          letter_spacing=None, wrap=True, fit=True):
    """Insert a text box with one paragraph + one run. Returns the shape.

    `fit` (default on) shrinks the font deterministically so the text fits
    its box in any viewer; `wrap` controls whether long text may break to a
    new line (set False for numbers/labels that must stay on one line).
    `auto_size` is also set so PowerPoint re-fits on open as a backstop."""
    tb = slide.shapes.add_textbox(in_(left), in_(top), in_(width), in_(height))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = wrap
    try:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except Exception:
        pass
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text.upper() if all_caps else text
    r.font.name = font
    eff = _fit_pt(r.text, float(width), float(height), float(size.pt),
                  wrap=wrap, all_caps=False) if fit else float(size.pt)
    r.font.size = Pt(eff)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    if letter_spacing is not None:
        # python-pptx doesn't expose spc directly via property; set on xml
        from pptx.oxml.ns import qn
        r._r.get_or_add_rPr().set("spc", str(letter_spacing))
    return tb


def _hrule(slide, left, top, width, color=MUTED, thick_pt=RULE_THICK_PT):
    """Thin horizontal rule. Implemented as a 0.02"-tall rect with fill."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  in_(left), in_(top), in_(width), in_(0.005))
    shp.line.fill.background()
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.shadow.inherit = False
    return shp


def _card(slide, left, top, width, height, *, fill=WHITE, border=MUTED,
          left_accent=GOLD, accent_w=LEFT_ACCENT_W_IN):
    """Bordered card with optional left-edge gold accent strip. Returns the
    main card shape (for layering text on top)."""
    # Main body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   in_(left), in_(top), in_(width), in_(height))
    body.fill.solid()
    body.fill.fore_color.rgb = fill
    body.line.color.rgb = border
    body.line.width = Pt(BORDER_THICK_PT)
    body.shadow.inherit = False
    # Left accent strip
    if left_accent is not None and accent_w > 0:
        accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                         in_(left), in_(top),
                                         in_(accent_w), in_(height))
        accent.fill.solid()
        accent.fill.fore_color.rgb = left_accent
        accent.line.fill.background()
        accent.shadow.inherit = False
    return body


def _metric_block(slide, left, top, width, label, value, *,
                  value_color=BLACK, value_size=SZ_VALUE):
    """Label-over-value metric primitive. label in muted 8.5pt caps,
    value in 14pt black.

    The value renders on a SINGLE line (wrap=False) so long figures like
    "HKD 581.4B" or a date "06 Aug 2026" shrink-to-fit instead of wrapping
    to two lines and spilling out of the row — a real overflow defect on
    high-cap / long-currency names."""
    _text(slide, left, top, width, 0.18, label, size=SZ_LABEL,
          color=MUTED, all_caps=True)
    _text(slide, left, top + 0.20, width, 0.32, value,
          size=value_size, color=value_color, bold=False, wrap=False)


def _section_label(slide, left, top, width, text):
    """10.5pt all-caps gray label with a horizontal rule above it."""
    _hrule(slide, left, top, width, color=MUTED, thick_pt=0.5)
    _text(slide, left, top + 0.10, width, 0.22, text,
          size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)


# ── Slide 1 sections ──────────────────────────────────────────

def _header_strip(slide, page_num: int, page_title: str):
    _text(slide, MARGIN_L, 0.32, 1.5, 0.28, "JABAL",
          font=FONT_DISPLAY, size=Pt(15), bold=True, color=BLACK)
    _text(slide, MARGIN_L, 0.58, 2.0, 0.18, "ASSET MANAGEMENT",
          size=Pt(8.5), color=GRAY, all_caps=True)
    _text(slide, 3.05, 0.36, 4.0, 0.22,
          f"PAGE {page_num}  ·  {page_title.upper()}",
          size=SZ_HEADER, color=BLACK, all_caps=True, align=PP_ALIGN.RIGHT)
    _text(slide, 3.05, 0.58, 4.0, 0.18, "INSTITUTIONAL RESEARCH  ·  EQUITY",
          size=Pt(8.5), color=MUTED, all_caps=True, align=PP_ALIGN.RIGHT)
    _hrule(slide, MARGIN_L, 0.88, CONTENT_W)


def _title_block(slide, company_name: str, meta_line: str, period: str):
    _text(slide, MARGIN_L, 1.08, CONTENT_W, 0.22, "EARNINGS PREVIEW NOTE",
          size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)
    _text(slide, MARGIN_L, 1.40, CONTENT_W, 0.62, company_name,
          font=FONT_DISPLAY, size=SZ_HERO, bold=False, color=BLACK)
    _text(slide, MARGIN_L, 2.08, CONTENT_W, 0.22, meta_line,
          size=SZ_META, color=GRAY)
    _text(slide, MARGIN_L, 2.36, CONTENT_W, 0.30, period,
          size=Pt(13), color=BLACK, bold=True)


def _consensus_row(slide, top: float, rating: str, n_analysts: int,
                    target_price: str, upside_pct: Optional[float]):
    _section_label(slide, MARGIN_L, top, CONTENT_W, "Analyst Consensus")
    row_top = top + 0.40
    card_w = (CONTENT_W - 0.22) / 3
    # Card 1: Rating
    _card(slide, MARGIN_L, row_top, card_w, 0.82, fill=CARD)
    _text(slide, MARGIN_L + 0.18, row_top + 0.08, card_w - 0.20, 0.18,
          f"RATING  ·  {n_analysts} ANALYSTS", size=SZ_LABEL, color=MUTED,
          all_caps=True)
    _text(slide, MARGIN_L + 0.18, row_top + 0.28, card_w - 0.20, 0.48,
          rating.upper(), size=SZ_VALUE_LG, color=BLACK, bold=True)
    # Card 2: Target Price
    c2_left = MARGIN_L + card_w + 0.11
    _card(slide, c2_left, row_top, card_w, 0.82, fill=CARD)
    _text(slide, c2_left + 0.18, row_top + 0.08, card_w - 0.20, 0.18,
          "TARGET PRICE", size=SZ_LABEL, color=MUTED, all_caps=True)
    _text(slide, c2_left + 0.18, row_top + 0.28, card_w - 0.20, 0.48,
          target_price, size=SZ_VALUE_LG, color=BLACK, bold=True)
    # Card 3: Upside
    c3_left = MARGIN_L + 2 * (card_w + 0.11)
    _card(slide, c3_left, row_top, card_w, 0.82, fill=CARD)
    # Label adapts to direction: a negative number under "UPSIDE TO
    # TARGET" misreads — call it DOWNSIDE / IMPLIED MOVE depending on sign.
    if upside_pct is None:
        label = "IMPLIED MOVE"
    elif upside_pct < 0:
        label = "DOWNSIDE TO TARGET"
    else:
        label = "UPSIDE TO TARGET"
    _text(slide, c3_left + 0.18, row_top + 0.08, card_w - 0.20, 0.18,
          label, size=SZ_LABEL, color=MUTED, all_caps=True)
    up_str = "—" if upside_pct is None else f"{upside_pct:+.1f}%"
    _text(slide, c3_left + 0.18, row_top + 0.28, card_w - 0.20, 0.48,
          up_str, size=SZ_VALUE_LG, color=signed_color(upside_pct),
          bold=True)


def _key_data_row(slide, top: float, items: list[tuple[str, str]]):
    """Six (label, value) tuples spread across 6.6"."""
    _section_label(slide, MARGIN_L, top, CONTENT_W, "Key Data")
    row_top = top + 0.40
    n = len(items)
    col_w = CONTENT_W / n
    for i, (label, value) in enumerate(items):
        left = MARGIN_L + i * col_w
        _metric_block(slide, left, row_top, col_w - 0.10, label, value)


def _performance_row(slide, top: float, items: list[tuple[str, Optional[float]]]):
    _section_label(slide, MARGIN_L, top, CONTENT_W, "Recent Performance")
    row_top = top + 0.40
    n = len(items)
    col_w = CONTENT_W / n
    for i, (label, pct) in enumerate(items):
        left = MARGIN_L + i * col_w
        val_str = "—" if pct is None else f"{pct:+.1f}%"
        _metric_block(slide, left, row_top, col_w - 0.10, label, val_str,
                       value_color=signed_color(pct))


def _range_bar(slide, top: float, low: float, high: float, current: float,
                currency: str):
    _section_label(slide, MARGIN_L, top, CONTENT_W, "52-Week Range")
    bar_top = top + 0.48
    # Low/high labels at ends
    _text(slide, MARGIN_L, bar_top - 0.04, 0.9, 0.20,
          f"{currency} {low:,.2f}", size=Pt(10), color=GRAY)
    _text(slide, MARGIN_L + CONTENT_W - 0.9, bar_top - 0.04, 0.9, 0.20,
          f"{currency} {high:,.2f}", size=Pt(10), color=GRAY,
          align=PP_ALIGN.RIGHT)
    # Track
    track_left = MARGIN_L + 0.95
    track_w = CONTENT_W - 1.90
    track = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    in_(track_left), in_(bar_top + 0.06),
                                    in_(track_w), in_(0.06))
    track.fill.solid()
    track.fill.fore_color.rgb = MUTED
    track.line.fill.background()
    # Fill from low to current
    if high > low:
        frac = max(0.0, min(1.0, (current - low) / (high - low)))
    else:
        frac = 0.5
    fill_w = track_w * frac
    fill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   in_(track_left), in_(bar_top + 0.06),
                                   in_(max(fill_w, 0.02)), in_(0.06))
    fill.fill.solid()
    fill.fill.fore_color.rgb = GOLD
    fill.line.fill.background()
    # Diamond marker at current
    marker_size = 0.18
    marker = slide.shapes.add_shape(MSO_SHAPE.DIAMOND,
                                     in_(track_left + fill_w - marker_size / 2),
                                     in_(bar_top + 0.06 - marker_size / 2 + 0.03),
                                     in_(marker_size), in_(marker_size))
    marker.fill.solid()
    marker.fill.fore_color.rgb = BLACK
    marker.line.fill.background()
    # Current label below marker
    _text(slide, track_left + fill_w - 0.7, bar_top + 0.24, 1.4, 0.20,
          f"Current  {currency} {current:,.2f}",
          size=Pt(9), color=BLACK, align=PP_ALIGN.CENTER)


def _highlights_row(slide, top: float, items: list[tuple[str, str]]):
    """List of (category_pill, body) rows. Up to 5."""
    _section_label(slide, MARGIN_L, top, CONTENT_W,
                    "Analyst Highlights  ·  Key Points")
    row_top = top + 0.50
    for i, (cat, body) in enumerate(items[:5]):
        y = row_top + i * 0.42
        pill = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                       in_(MARGIN_L), in_(y),
                                       in_(0.85), in_(0.22))
        pill.fill.solid()
        pill.fill.fore_color.rgb = CARD
        pill.line.fill.background()
        _text(slide, MARGIN_L, y, 0.85, 0.22, cat,
              size=SZ_BULLET_PILL, color=GRAY, all_caps=True, bold=True,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        _text(slide, MARGIN_L + 0.98, y - 0.02, CONTENT_W - 0.98, 0.32, body,
              size=SZ_BODY, color=BLACK)


def _footer(slide, page_num: int, total_pages: int, sources: str,
             analyst_name: str, gen_date: str, freshness_banner: str = ""):
    top = 12.47
    # When a freshness banner is provided, drop it just above the footer
    # rule so the analyst sees data-age at a glance without opening the
    # provenance xlsx.
    if freshness_banner:
        _text(slide, MARGIN_L, top - 0.20, CONTENT_W, 0.18,
              freshness_banner,
              size=SZ_FOOTER, color=GRAY, italic=True)
    _hrule(slide, MARGIN_L, top, CONTENT_W)
    _text(slide, MARGIN_L, top + 0.07, CONTENT_W, 0.18,
          f"Source: {sources}  |  Generated {gen_date}  |  Analyst: {analyst_name}",
          size=SZ_FOOTER, color=GRAY)
    _text(slide, MARGIN_L, 12.90, 5.5, 0.18,
          "Jabal Asset Management  ·  Regulated by the Financial Services Authority of Oman",
          size=SZ_FOOTER, color=GRAY)
    _text(slide, 6.05, 12.90, 1.0, 0.18,
          f"{page_num} / {total_pages}",
          size=SZ_TAB_NUM, color=GRAY, align=PP_ALIGN.RIGHT)
    _text(slide, MARGIN_L, 13.08, CONTENT_W, 0.16,
          "CONFIDENTIAL  ·  For Institutional & Qualified Investors Only",
          size=Pt(7.5), color=MUTED, all_caps=True, align=PP_ALIGN.CENTER)


# ── Public entry point ────────────────────────────────────────

@dataclass
class SnapshotData:
    """The exact set of inputs Slide 1 needs. Built by the orchestrator
    from canonical_store + report metadata."""
    company_name: str
    ticker: str
    sector: str
    industry: str
    exchange: str
    period_label: str               # "Q2 2026 Earnings Preview"
    rating: str                     # "OUTPERFORM"
    n_analysts: int
    target_price_fmt: str           # "SAR 137.70"
    upside_pct: Optional[float]     # -4.5
    last_close_fmt: str
    market_cap_fmt: str
    report_date: str
    pe_fy_est_fmt: str
    div_yield_fmt: str
    currency: str
    perf_1d: Optional[float]
    perf_1w: Optional[float]
    perf_1m: Optional[float]
    perf_3m: Optional[float]
    perf_6m: Optional[float]
    perf_ytd: Optional[float]
    range_low: float
    range_high: float
    range_current: float
    highlights: list[tuple[str, str]]   # (CATEGORY, body) — max 5
    sources_line: str
    analyst_name: str
    gen_date: str
    pe_fy_est_label: str = "P/E (FY EST)"  # e.g. "P/E (FY26E)" when year known
    total_pages: int = 3
    # Data-freshness banner — one short line above the footer that tells
    # the analyst at a glance how recent each tier of data is. Set by
    # build_snapshot_data based on canonical_store cell timestamps +
    # the live-quote record (if applied).
    freshness_banner: str = ""
    # Numeric price the deck displayed as LAST CLOSE (and computed UPSIDE
    # from). The provenance reuses this so the .xlsx "Last close" and
    # "Upside to target" trace to the SAME number the slide shows — not a
    # separately re-read canonical price that may have drifted.
    last_close_value: Optional[float] = None


def render_snapshot_slide(prs, data: SnapshotData):
    """Add slide 1 to an existing python-pptx Presentation. Returns the slide."""
    layout = prs.slide_layouts[6] if len(prs.slide_layouts) > 6 else prs.slide_layouts[-1]
    # Prefer a blank layout
    blank = next((L for L in prs.slide_layouts if L.name.lower() == "blank"),
                  layout)
    slide = prs.slides.add_slide(blank)

    # Sections
    _header_strip(slide, 1, "Snapshot")
    _title_block(
        slide, data.company_name,
        f"{data.ticker}  ·  {data.sector}  ·  {data.industry}  ·  {data.exchange}",
        data.period_label,
    )
    _consensus_row(slide, 2.96, data.rating, data.n_analysts,
                    data.target_price_fmt, data.upside_pct)
    _key_data_row(slide, 4.38, [
        ("LAST CLOSE", data.last_close_fmt),
        ("MARKET CAP", data.market_cap_fmt),
        # The date the company will NEXT report — not the deck's generation
        # date (footer). "EARNINGS DATE" was still read as ambiguous, so use
        # "NEXT EARNINGS" which can only mean the upcoming print.
        ("NEXT EARNINGS", data.report_date),
        (data.pe_fy_est_label or "P/E (FY EST)", data.pe_fy_est_fmt),
        # Dividend yield is current — quote_price × shares / last paid
        # 12-month total — so label it "TTM" to disambiguate from the
        # forward yield implied by FY-est dividend forecasts.
        ("DIV YIELD (TTM)", data.div_yield_fmt),
        ("CURRENCY", data.currency),
    ])
    _performance_row(slide, 5.42, [
        ("1 DAY", data.perf_1d),
        ("1 WEEK", data.perf_1w),
        ("1 MONTH", data.perf_1m),
        ("3 MONTHS", data.perf_3m),
        ("6 MONTHS", data.perf_6m),
        ("YTD", data.perf_ytd),
    ])
    _range_bar(slide, 6.46, data.range_low, data.range_high,
                data.range_current, data.currency)
    _highlights_row(slide, 7.62, data.highlights)
    _footer(slide, 1, data.total_pages, data.sources_line,
             data.analyst_name, data.gen_date,
             freshness_banner=getattr(data, "freshness_banner", ""))
    return slide


# ── Rating-label prettifier ───────────────────────────────────

def _pretty_rating(raw) -> str:
    """Normalise consensus-rating strings from any provider into the
    title-case form analysts use in print: e.g.
        'STRONG_BUY' -> 'Strong Buy'
        'OUTPERFORM' -> 'Outperform'
        'buy'        -> 'Buy'
        'hold/maintain' -> 'Hold/Maintain'
    Returns '' on empty/None input. Underscores and lower-case bleed
    through from Investing.com's enum (consensus_recommendation).
    """
    if raw is None:
        return ""
    s = str(raw).strip()
    if not s:
        return ""
    s = s.replace("_", " ").replace("-", " ")
    # Title-case but preserve runs of letters as words
    return " ".join(part.capitalize() for part in s.split())


# ── Sources-line builder ──────────────────────────────────────

def _sources_line(cv: dict) -> str:
    """Build the slide-footer "Source:" line.

    Lists every provider that contributed a value to ANY canonical field
    (i.e. the union of `sources_with_value`), not just the per-field
    winners. This means Investing.com / Yahoo / MarketScreener all show
    up when they fed the deck, even when only one of them won the
    reconciliation for a given cell.
    """
    seen: set[str] = set()
    for c in cv.values():
        winner = getattr(c, "canonical_source", "") or ""
        if winner:
            seen.add(winner)
        for s in (getattr(c, "sources_with_value", None) or []):
            if s:
                seen.add(s)
    # Primary → Secondary order: Bloomberg (when uploaded) > Investing.com >
    # Yahoo Finance > MarketScreener > supporting layers. MS is always
    # demoted because it's web-scraped (vs API-fed) and its absolute values
    # have shown systematic mismatches (e.g. H-share-only market cap for
    # HK names) vs the live Investing equity-page numbers.
    pretty = {
        "bloomberg":      "Bloomberg",
        "ir_pdf":         "Company IR",
        "investing":      "Investing.com",
        "yahoo":          "Yahoo Finance",
        "marketscreener": "MarketScreener",
        "macro":          "World Bank / IMF",
        "ishares":        "iShares",
        "commodities":    "World Bank / OPEC / EIA",
    }
    _order = {
        "bloomberg":      0,
        "ir_pdf":         1,
        "investing":      2,
        "yahoo":          3,
        "marketscreener": 4,
        "macro":          5,
        "ishares":        6,
        "commodities":    7,
    }
    ordered = sorted(seen, key=lambda s: (_order.get(s, 99), s))
    labels = [pretty.get(s, s) for s in ordered]
    return ", ".join(labels) or "free-source stack"


# ── Highlight derivation ──────────────────────────────────────

def _llm_highlights_or_template(*, ticker: str, cv: dict, currency: str,
                                   current_price, mcap, target_mean, upside_pct,
                                   pe_fwd, div_yield, range_low, range_high,
                                   n_analysts: int, rs: dict) -> list[tuple[str, str]]:
    """Prefer Gemini-generated highlights when available; fall back to the
    deterministic template. The LLM bodies have already been numeric-trace
    validated upstream — any rejected category falls back to its template
    line so the slide always shows 5 pills.

    Categories must appear in this order: EARNINGS, VALUATION, POSITIONING,
    WATCH, RISK. The template line is always computed so we can backfill
    any slot the LLM dropped.
    """
    template = _derive_highlights(
        cv=cv, currency=currency, current_price=current_price,
        mcap=mcap, target_mean=target_mean, upside_pct=upside_pct,
        pe_fwd=pe_fwd, div_yield=div_yield,
        range_low=range_low, range_high=range_high,
        n_analysts=n_analysts, rs=rs,
    )
    try:
        from src.services.llm_summary import generate_summary
        llm = generate_summary(ticker)
    except Exception:
        llm = None
    if not llm or not llm.get("highlights"):
        return template
    by_cat = {(item.get("category") or "").strip().upper(): (item.get("body") or "").strip()
              for item in (llm.get("highlights") or []) if isinstance(item, dict)}
    out: list[tuple[str, str]] = []
    cats = ("EARNINGS", "VALUATION", "POSITIONING", "WATCH", "RISK")
    # Index the template by its category label for quick fallback lookup.
    tmpl_by_cat = {cat.upper(): body for cat, body in template}
    for cat in cats:
        body = by_cat.get(cat) or tmpl_by_cat.get(cat, "")
        if body:
            out.append((cat, body))
    return out or template


def _derive_highlights(*, cv: dict, currency: str, current_price,
                         mcap, target_mean, upside_pct, pe_fwd,
                         div_yield, range_low, range_high,
                         n_analysts: int, rs: dict) -> list[tuple[str, str]]:
    """Derive 5 highlight pills from real numbers in canonical_store.
    Each row is anchored to a specific figure that also appears in the
    rest of the deck — keeps the slide internally consistent and avoids
    hardcoded boilerplate. Falls back to a short generic line only when
    no numeric anchor is available."""
    cur = currency or ""
    rows: list[tuple[str, str]] = []

    # EARNINGS — anchor on consensus rating + analyst count, with a fallback
    # to the Investing surprise track when present.
    rating_label = ""
    if isinstance(rs, dict):
        rating_label = _pretty_rating(rs.get("consensus"))
    if rating_label and n_analysts:
        rows.append(("EARNINGS",
            f"A {rating_label} consensus across {n_analysts} analysts sets a high bar the print must clear to hold the multiple."))
    elif n_analysts:
        rows.append(("EARNINGS",
            f"With {n_analysts} analysts covering, the print must validate the earnings trajectory already priced in."))
    else:
        rows.append(("EARNINGS",
            "The print tests whether the recent earnings trajectory can be sustained into the next quarter."))

    # VALUATION — P/E (FY est) is the most universally available figure.
    if isinstance(pe_fwd, (int, float)) and pe_fwd > 0:
        rows.append(("VALUATION", f"Forward P/E {float(pe_fwd):.1f}x — anchor for re-rate / de-rate debate."))
    elif isinstance(mcap, (int, float)) and mcap > 0:
        # No P/E available — fall back to market-cap framing.
        if mcap >= 1e12:
            rows.append(("VALUATION", f"Market cap {cur} {mcap/1e12:.2f}T — index-eligible scale."))
        elif mcap >= 1e9:
            rows.append(("VALUATION", f"Market cap {cur} {mcap/1e9:.1f}B — institutional-grade liquidity."))
        else:
            rows.append(("VALUATION", "Valuation context limited; awaiting MS / Investing refresh."))
    else:
        rows.append(("VALUATION", "Valuation context limited; awaiting MS / Investing refresh."))

    # POSITIONING — dividend yield is the strongest single signal for GCC banks,
    # SOEs, and dividend-heavy names. Fall back to current vs 52-week high
    # for growth names with no dividend.
    if isinstance(div_yield, (int, float)) and div_yield > 0:
        rows.append(("POSITIONING",
            f"The {float(div_yield):.2f}% yield anchors income holders, but that base also limits appetite to chase a soft print."))
    elif (isinstance(current_price, (int, float)) and current_price > 0
          and isinstance(range_high, (int, float)) and range_high > 0):
        gap = (current_price / range_high - 1.0) * 100
        rows.append(("POSITIONING",
            f"Trading {gap:+.1f}% off the 52-week high — the debate is whether the pullback is value or a warning."))
    else:
        rows.append(("POSITIONING",
            "Positioning hinges on whether the print breaks the recent trading range either way."))

    # WATCH — the swing factor heading into the print. Fall back to a
    # mechanism (guidance closing the target gap), never a bare datapoint.
    if (isinstance(upside_pct, (int, float))
        and isinstance(target_mean, (int, float)) and target_mean > 0):
        rows.append(("WATCH",
            f"Whether guidance confirms the trajectory needed to close the {upside_pct:+.1f}% gap to the {cur} {target_mean:,.2f} target."))
    elif isinstance(target_mean, (int, float)) and target_mean > 0:
        rows.append(("WATCH",
            "Management's forward guidance is the swing factor for whether the consensus target holds."))
    else:
        rows.append(("WATCH", "Management commentary on the forward outlook is the swing factor for the print."))

    # RISK — analyst-distribution concentration is the most defensible
    # quantitative risk anchor (one-sided consensus = harder to surprise).
    # Some providers (notably MS for GCC names) publish only the consensus
    # label + analyst total without a buy/hold/sell breakdown; infer the
    # skew from the consensus label in that case so the pill still anchors
    # on a real signal.
    if isinstance(rs, dict):
        buy = int(rs.get("buy", 0) or 0)
        hold = int(rs.get("hold", 0) or 0)
        sell = int(rs.get("sell", 0) or 0)
        total = buy + hold + sell
        if total == 0 and n_analysts > 0:
            raw = (rs.get("consensus") or "").upper()
            pretty = _pretty_rating(rs.get("consensus"))
            if any(t in raw for t in ("BUY", "OUTPERFORM", "ACCUMULATE")):
                rows.append(("RISK",
                    f"Crowded long — consensus {pretty} across {n_analysts} analysts raises the expectations bar."))
            elif any(t in raw for t in ("SELL", "UNDERPERFORM", "REDUCE")):
                rows.append(("RISK",
                    f"Tape skewed bearish — consensus {pretty} across {n_analysts} analysts."))
            elif pretty:
                rows.append(("RISK",
                    f"View dispersion — consensus {pretty} across {n_analysts} analysts."))
            else:
                rows.append(("RISK", f"{n_analysts} analysts covering; breakdown not disclosed."))
        elif total > 0:
            denom = max(1, total)
            buy_share  = buy  / denom
            hold_share = hold / denom
            sell_share = sell / denom
            # Order matters: detect dominant Hold first (60%+ Hold isn't
            # "one-sided Buy" even when sells are 0).
            if hold_share >= 0.50 and total >= 5:
                rows.append(("RISK",
                    f"Consensus mixed — {hold}/{total} ratings are Hold despite Buy screen ({buy} Buy · {sell} Sell)."))
            elif buy_share >= 0.80 and total >= 5:
                rows.append(("RISK",
                    f"Crowded long — {buy}/{total} buy ratings raise the expectations bar."))
            elif sell_share >= 0.30 and total >= 5:
                rows.append(("RISK",
                    f"Tape skewed bearish — {sell}/{total} sell ratings ({sell_share*100:.0f}%)."))
            elif sell == 0 and buy_share >= 0.60 and total >= 5:
                rows.append(("RISK",
                    f"No bears in the tape — {buy}/{total} buys, 0 sells across {total} analysts."))
            else:
                rows.append(("RISK",
                    "A print short of consensus would trigger estimate cuts and a de-rating of the multiple."))
        else:
            rows.append(("RISK",
                "A downside surprise on the print would pressure both forward estimates and the multiple."))
    else:
        rows.append(("RISK", "Macro / sector sensitivity; refer to thesis on slide 2."))

    return rows


# ── Data adapter: canonical_store → SnapshotData ──────────────

def _perf_and_range_from_series(close_series, live_price) -> dict:
    """Compute recent-performance %s and the 52-week range the way the
    Investing/MarketScreener WEBSITE does: the live price vs the close N
    calendar days back, and the min/max of a full-year close series.

    WHY: Investing's historical-data *API* returns perf fields that are
    pre-computed and lag its live quote (its `perf_updated_at` trails the
    site by days), so the deck's 1D/1W/… diverged from what the analyst sees
    on the site. Anchoring on the live price reproduces the site's numbers
    (verified: 1D/1W/1M/3M match to the basis point on BKMB.OM). The 52-week
    range is taken from the series min/max ONLY when the series actually
    spans ~a year — a truncated series can't fabricate a wrong range
    (the "0.37–0.45 instead of 0.27–0.49" bug)."""
    out: dict = {}
    if not (isinstance(close_series, list) and len(close_series) > 5
            and isinstance(live_price, (int, float)) and live_price > 0):
        return out
    from datetime import datetime as _dt, timedelta as _td, date as _date
    series = []
    for p in close_series:
        try:
            series.append((_dt.strptime(p["date"], "%Y-%m-%d").date(), float(p["close"])))
        except (KeyError, TypeError, ValueError):
            continue
    if len(series) < 5:
        return out
    series.sort()
    last_date, last_close = series[-1]
    span_days = (last_date - series[0][0]).days

    # PERF anchor (numerator) = the series' OWN last close. This makes the
    # performance the series' internal ratio (last_close / close_N_ago), which
    # is exactly what the source website shows — and is INDEPENDENT of the
    # (possibly stale or differently-sourced) canonical price used for display.
    # BKMB's canonical price was 0.41 while its real latest close was 0.400;
    # anchoring on the series fixes the perf without touching the display.
    anchor = last_close
    out["latest_close"] = last_close

    # Granularity: only RECOMPUTE perf when the series is daily. Yahoo stores a
    # DOWNSAMPLED WEEKLY close-series (but computes its own accurate daily
    # perf_*), so recomputing 1D/1W from weekly points grabs a week-old anchor
    # and is wrong (9988.HK 1D came out -4.6% vs the real -1.4%). For coarse
    # series we leave perf to the source's pre-computed fields and only take
    # the 52-week range from the span. Investing's series is daily → recompute
    # (its pre-computed perf lags, the BKMB case).
    gaps = sorted((series[i + 1][0] - series[i][0]).days for i in range(len(series) - 1))
    median_gap = gaps[len(gaps) // 2] if gaps else 99
    is_daily = median_gap <= 2

    def _base(n):
        target = last_date - _td(days=n)
        # Reject an anchor that is far older than the target (coarse series):
        cand = [(d, c) for d, c in series if d <= target]
        if not cand:
            return None
        bd, bc = cand[-1]
        return bc if (target - bd).days <= max(4, n * 0.5) else None

    if is_daily:
        for key, n in (("perf_1d", 1), ("perf_1w", 7), ("perf_1m", 30),
                        ("perf_3m", 91), ("perf_6m", 182)):
            if span_days >= n:
                b = _base(n)
                if b:
                    out[key] = round((anchor / b - 1.0) * 100, 2)
        jan1 = _date(last_date.year, 1, 1)
        ytd = [c for d, c in series if d <= jan1]
        if ytd and span_days >= (last_date - jan1).days:
            out["perf_ytd"] = round((anchor / ytd[-1] - 1.0) * 100, 2)
    if span_days >= 330:                       # trust the range only on a full year
        out["range_52w_low"] = min(c for _, c in series)
        out["range_52w_high"] = max(c for _, c in series)
    out["_span_days"] = span_days
    return out


def _price_asof_from_history(hist_prices: dict | None) -> "datetime | None":
    """The true as-of date of the price = the last dated close in the
    historical series. This is dated source data we can trust, unlike the
    canonical store's last_refreshed_at (= read time). Returns an aware
    UTC datetime, or None when no dated series is available."""
    if not isinstance(hist_prices, dict):
        return None
    series = hist_prices.get("close_series") or []
    if not isinstance(series, list) or not series:
        return None
    from datetime import datetime as _dt, timezone as _tz
    last_date = None
    for pt in series:
        d = pt.get("date") if isinstance(pt, dict) else None
        if not d:
            continue
        try:
            parsed = _dt.fromisoformat(str(d)[:10]).replace(tzinfo=_tz.utc)
        except Exception:
            continue
        if last_date is None or parsed > last_date:
            last_date = parsed
    return last_date


def _compute_freshness_banner(cv: dict, live_quote_record: dict | None,
                                price_asof: "datetime | None" = None) -> str:
    """Build the one-line freshness summary that drops above the slide-1
    footer. Tells the analyst at a glance which fields are intraday-live,
    which are daily-snapshot, and which are macro-cadence — so they
    never wonder how stale the deck is.

    Tiers we surface:
      * Live (yfinance, < 15 min)  — when live_quote_record is present
      * Daily snapshot age         — TRUE age of the price data
      * Macro                       — IMF/WB reporting year (always slow)

    `price_asof` is the actual as-of date of the price data (the last
    close in the historical series). It MUST be used in preference to the
    canonical store's `last_refreshed_at`, which records when the pipeline
    *read* the (possibly cached/stale) snapshot — not when the underlying
    data is from. Using last_refreshed_at made the banner claim "0h old"
    on a 16-day-old Render snapshot; the deck's own rule is "data ≤1 day
    old or the deck says so", so the age shown here must be the real one.
    """
    from datetime import datetime as _dtf, timezone as _tzf

    def _age_phrase(asof) -> str:
        """Render a true age as 'Nh old' / 'Nd old' from an aware datetime."""
        try:
            if asof.tzinfo is None:
                asof = asof.replace(tzinfo=_tzf.utc)
            age_h = int((_dtf.now(_tzf.utc) - asof).total_seconds() / 3600)
            age_h = max(0, age_h)
            if age_h < 24:
                return f"Price snapshot {age_h}h old"
            return f"Price snapshot {age_h // 24}d old"
        except Exception:
            return ""

    parts: list[str] = []
    if live_quote_record and live_quote_record.get("ok"):
        fetched = live_quote_record.get("fetched_at", "")
        try:
            t_fetched = _dtf.fromisoformat(fetched.replace("Z", "+00:00"))
            now = _dtf.now(_tzf.utc)
            mins = max(0, int((now - t_fetched).total_seconds() / 60))
            parts.append(f"Price live ({mins}m ago)")
        except Exception:
            parts.append("Price live")
    else:
        # No live quote — surface the TRUE snapshot age. Prefer the price
        # data's real as-of date (last close in the historical series);
        # fall back to last_refreshed_at only when no dated price is
        # available (which at least won't *under*-state age in practice).
        phrase = ""
        if price_asof is not None:
            phrase = _age_phrase(price_asof)
        if not phrase:
            cp = cv.get("current_price")
            if cp and cp.last_refreshed_at:
                phrase = _age_phrase(cp.last_refreshed_at)
        if phrase:
            parts.append(phrase)

    # Daily-snapshot tier (target, dividend yield, rating split).
    snapshot_ages_h: list[int] = []
    from datetime import datetime as _dtf2, timezone as _tzf2
    def _has_real_value(c) -> bool:
        """A forecast cell only counts toward freshness when it actually
        carries data. Keying on the timestamp alone made the banner claim
        'Forecasts < 1d old' for a stale cell whose value had failed to
        fetch (empty dict / None) — a freshness lie. Require a real value."""
        if not c:
            return False
        v = getattr(c, "value", None)
        if v is None:
            return False
        if isinstance(v, dict):
            return any(x is not None for x in v.values())
        if isinstance(v, (list, tuple, str)):
            return len(v) > 0
        return True

    for field in ("target_price", "dividend_yield", "rating_split",
                   "valuation_forward", "valuation_historical"):
        c = cv.get(field)
        if c and c.last_refreshed_at and _has_real_value(c):
            try:
                snapshot_ages_h.append(int(
                    (_dtf2.now(_tzf2.utc) - c.last_refreshed_at).total_seconds() / 3600))
            except Exception:
                pass
    if snapshot_ages_h:
        max_age_h = max(snapshot_ages_h)
        if max_age_h < 24:
            parts.append(f"Forecasts < 1d old")
        elif max_age_h < 48:
            parts.append(f"Forecasts ~1d old")
        else:
            parts.append(f"Forecasts {max_age_h // 24}d old")

    # Macro tier — always slow, IMF/WB.
    prof = cv.get("company_profile")
    macro_year = ""
    if prof and isinstance(prof.value, dict):
        macro_year = str(prof.value.get("macro_year")
                          or prof.value.get("gdp_growth_fcst_year") or "")
    if macro_year:
        parts.append(f"Macro IMF {macro_year}")

    if not parts:
        return ""
    return "Data freshness:  " + "  ·  ".join(parts)


def build_snapshot_data(ticker: str, *, analyst_name: str = "Jabal Research",
                          period_label: str = "Q2 2026 Earnings Preview",
                          report_date: str = "TBA",
                          highlights: Optional[list[tuple[str, str]]] = None,
                          ms_price_performance: Optional[dict] = None,
                          historical_override: Optional[dict] = None,
                          live_quote_record: Optional[dict] = None,
                          ) -> SnapshotData:
    """Translate canonical_store rows into the slide's input dataclass.
    Defensive against missing fields: every renderer-visible string has
    a sensible default ('—' for numerics, '' for text)."""
    cv = get_all_fields(ticker)

    def _val(field):
        c = cv.get(field)
        return c.value if c else None

    profile = _val("company_profile") or {}
    if not isinstance(profile, dict):
        profile = {}

    # Backfill profile from company_master when canonical_store hasn't been
    # populated by a recent yfinance probe (common on first-run tickers and
    # after a Render restart that wiped /tmp). This gives us at minimum a
    # sector / industry / exchange / country / company name to render in the
    # header subtitle — preventing the "· —" tail seen on fresh runs.
    try:
        from src.storage.db import load_company as _load_company
        cm = _load_company(ticker) or {}
        if cm:
            profile.setdefault("name", cm.get("company_name") or "")
            profile.setdefault("sector", cm.get("sector") or "")
            profile.setdefault("industry", cm.get("industry") or "")
            # Currency: company_master overrides — some providers report a
            # company's country-of-record currency (Tencent: CNY) but the
            # stock trades in the listing currency (HKD). The deck should
            # always show the trading currency.
            if cm.get("currency"):
                profile["currency"] = cm["currency"]
            # Exchange suffix: friendly-name + country, e.g. "Tadawul (Saudi Arabia)".
            # The DB stores 3-letter codes (SAU, ADX, NSE, HKG, ...) — map to
            # human-readable bourse names. Unknown codes fall through to the raw code.
            _EX_NAMES = {
                "SAU": "Tadawul", "ADX": "ADX", "DFM": "DFM",
                "MSM": "MSX",  "DSM": "QSE", "BHB": "Bahrain Bourse",
                "KSE": "Boursa Kuwait", "EGX": "EGX",
                "NSE": "NSE",  "BSE": "BSE",
                "HKG": "HKEX", "SHA": "SSE", "SHZ": "SZSE",
                "TYO": "TSE",  "KRX": "KRX",
                "JNB": "JSE",
                "NMS": "NASDAQ", "NCM": "NASDAQ", "NGM": "NASDAQ",
                "NYQ": "NYSE", "NYS": "NYSE", "ASE": "NYSE American",
                "LON": "LSE", "PAR": "Euronext Paris",
                "AMS": "Euronext Amsterdam", "BRU": "Euronext Brussels",
                "FRA": "Frankfurt", "STO": "Nasdaq Stockholm",
                "ASX": "ASX", "TSE": "TSX",
            }
            xcode = (cm.get("exchange") or "").strip().upper()
            country = (cm.get("country") or "").strip()
            if xcode and country:
                profile["exchange"] = f"{_EX_NAMES.get(xcode, xcode)} ({country})"
            elif xcode:
                profile["exchange"] = _EX_NAMES.get(xcode, xcode)
    except Exception:
        pass
    last_price = _val("current_price")
    mcap = _val("market_cap")
    val_hist = _val("valuation_historical") or {}
    val_fwd = _val("valuation_forward") or {}
    target = _val("target_price")
    rating_split = _val("rating_split") or {}
    div_yield = _val("dividend_yield")
    hist_prices = _val("historical_prices") or {}
    # Prefer the Investing override whenever canonical lacks a REAL close
    # series. For GCC ex-Saudi names canonical_store often holds only an
    # iShares-EEM proxy (ETF metadata, no close_series/perf), which would
    # otherwise make the deck fall through to MarketScreener's performance
    # numbers — disagreeing with both Investing.com and our own provenance.
    # Checking close_series (not mere truthiness) routes price/52w/perf
    # through the Investing series the analyst actually cross-checks against.
    if not (isinstance(hist_prices, dict) and hist_prices.get("close_series")) \
       and isinstance(historical_override, dict) and historical_override.get("close_series"):
        hist_prices = historical_override

    # Currency: try profile first, else canonical_value units
    currency = (profile.get("currency") if isinstance(profile, dict) else None) or ""
    if not currency and cv.get("current_price"):
        currency = cv["current_price"].canonical_source[:3].upper()  # fallback

    # Sub-unit currency normalisation. The South African JSE quotes
    # in ZAc (cents — 1/100 ZAR); registry carries a currency_unit_scale
    # of 100 so the displayed price reads as Rand rather than cents.
    # Other markets quote in base unit so scale=1 is a no-op.
    try:
        from src.services.ticker_registry import get_ticker_info as _gti
        _ts = _gti(ticker)
        unit_scale = float(_ts.get("currency_unit_scale") or 1)
        # When scale > 1, divide displayed numerics and relabel currency
        # to the base unit (ZAc → ZAR, GBp → GBP, ILA → ILS).
        if unit_scale and unit_scale != 1:
            _BASE_CCY = {"ZAc": "ZAR", "GBp": "GBP", "ILA": "ILS", "ILa": "ILS"}
            currency = _BASE_CCY.get(currency, currency)
    except Exception:
        unit_scale = 1

    def _scale(x):
        if not isinstance(x, (int, float)): return x
        try:
            return float(x) / unit_scale
        except (TypeError, ZeroDivisionError):
            return x

    # Format helpers
    def _price_dp(v: float) -> int:
        """Decimals scaled to price magnitude. Sub-1.0 prices (e.g. OMR names
        ~0.45) need 3+ dp so the displayed last close, target, and the
        upside% reconcile — at 2 dp a 0.5368 target shows '0.54' and the
        upside (computed from the unrounded value) looks wrong vs the
        displayed numbers."""
        av = abs(v)
        if av == 0 or av >= 1:
            return 2
        if av >= 0.1:
            return 3
        if av >= 0.01:
            return 4
        return 5

    def _money(x):
        if x is None:
            return "—"
        try:
            v = _scale(float(x))
            dp = _price_dp(v)
            return f"{currency} {v:,.{dp}f}" if currency else f"{v:,.{dp}f}"
        except (TypeError, ValueError):
            return "—"

    # Different sources report market_cap in different units:
    #   MarketScreener: millions of local currency
    #   Yahoo / ADX / HKEX / NSE: raw local currency units
    mc_source = cv.get("market_cap").canonical_source if cv.get("market_cap") else ""
    mc_scale = 1_000_000.0 if mc_source == "marketscreener" else 1.0

    def _mc(x):
        if x is None:
            return "—"
        try:
            v = float(x) * mc_scale
        except (TypeError, ValueError):
            return "—"
        if v >= 1e12:
            return f"{currency} {v/1e12:.2f}T" if currency else f"{v/1e12:.2f}T"
        if v >= 1e9:
            return f"{currency} {v/1e9:.1f}B" if currency else f"{v/1e9:.1f}B"
        if v >= 1e6:
            return f"{currency} {v/1e6:.0f}M" if currency else f"{v/1e6:.0f}M"
        return f"{currency} {v:,.0f}"

    def _pct(x):
        if x is None:
            return "—"
        try:
            return f"{float(x):.2f}%"
        except (TypeError, ValueError):
            return "—"

    # Pull rating + analyst count
    rating = "—"
    n_analysts = 0
    if isinstance(rating_split, dict):
        # MS shape varies: ideally {"buy":6,"hold":3,"sell":1,"total":10,"consensus":"OUTPERFORM"}
        # but sometimes only {"consensus":"OUTPERFORM"}.
        n_analysts = int(rating_split.get("total", 0) or 0)
        rating = _pretty_rating(rating_split.get("consensus")) or "—"
    # Fall back to target_price.n_analysts if rating_split didn't carry the total.
    if not n_analysts and isinstance(target, dict):
        n_analysts = int(target.get("n_analysts", 0) or 0)

    # P/E forward — first try MS's historical P/E series picked at the
    # current fiscal year. Previously we read `reversed()` and took the
    # last entry, which is the FURTHEST forecast year (e.g. FY2028 for a
    # FY24-28 series). The slide-1 chip should reflect the next reporting
    # FY, not the multi-year-out estimate. Falls through to valuation_forward
    # (Investing-derived pe_fy1) when historical isn't loaded.
    pe_fwd = None
    pe_fwd_year = None
    if isinstance(val_hist, dict):
        periods = val_hist.get("periods", []) or []
        pe_vals = val_hist.get("pe", []) or []
        if len(periods) == len(pe_vals):
            from datetime import datetime as _dt
            cur_year = _dt.now().year
            # Pick the entry whose period label contains the current year
            # (FY2026), or the closest next year if none matches.
            best_idx, best_diff, best_year = None, None, None
            for i, p in enumerate(periods):
                import re as _r
                m = _r.search(r"(\d{4})", str(p) or "")
                if not m: continue
                yr = int(m.group(1))
                if not isinstance(pe_vals[i], (int, float)): continue
                diff = abs(yr - cur_year)
                if best_diff is None or diff < best_diff:
                    best_idx, best_diff, best_year = i, diff, yr
            if best_idx is not None:
                pe_fwd = float(pe_vals[best_idx])
                pe_fwd_year = best_year
    if pe_fwd is None and isinstance(val_fwd, dict):
        from datetime import datetime as _dt
        _cy = _dt.now().year
        # Yahoo's valuation_forward bundle uses `forward_pe` (the FY+1 anchor);
        # Investing/MS use `pe_fy1` / `pe_2026`. Check Yahoo's key first so the
        # slide-1 chip fills for every Yahoo-covered name.
        for k in ("forward_pe", "pe_fy1", "pe_2026", "pe_2027", "pe"):
            v = val_fwd.get(k)
            if isinstance(v, (int, float)) and v > 0:
                pe_fwd = v
                # forward_pe / pe_fy1 = current fiscal year; pe_2026/2027 explicit.
                import re as _r
                m = _r.search(r"(\d{4})", k)
                pe_fwd_year = int(m.group(1)) if m else _cy
                break
    elif pe_fwd is None and isinstance(val_fwd, (int, float)):
        pe_fwd = val_fwd

    # Final fallback for slide-1 P/E (FY EST): synthesize from MS annual
    # net_income + market_cap when /valuation/ and Investing both empty.
    # Mirrors the synth path in render_jabal_valuation.py — see that
    # file for the math. Without this the BKMB slide-1 chip shows '—'
    # even though the canonical_store has everything we need.
    if pe_fwd is None and isinstance(mcap, (int, float)) and mcap > 0 \
       and isinstance(last_price, (int, float)) and last_price > 0:
        try:
            from src.storage.db import get_conn as _gc_pe
            import json as _json_pe
            _conn = _gc_pe()
            _row = _conn.execute(
                "SELECT payload_json FROM raw_observations "
                "WHERE provider='marketscreener' AND ticker=? AND field='income_statement_annual' "
                "ORDER BY rowid DESC LIMIT 1",
                (ticker,),
            ).fetchone()
            _conn.close()
            _ms_ann = (_json_pe.loads(_row["payload_json"]) if _row else {}).get("annual") or {}
        except Exception:
            _ms_ann = {}
        # Apply the same MS scale (1e6) the snapshot uses for market cap.
        _mc_source = cv.get("market_cap").canonical_source if cv.get("market_cap") else ""
        _mcap_raw = float(mcap) * (1_000_000.0 if (_mc_source or "").lower() == "marketscreener" else 1.0)
        _shares = _mcap_raw / float(last_price) if last_price else None
        _ni_series = _ms_ann.get("net_income") or []
        _periods_raw = _ms_ann.get("periods") or []
        from datetime import datetime as _dt2
        _cur_yr = _dt2.now().year
        if _shares and _ni_series:
            # Pick the NI value for the current calendar year (or closest).
            _best_idx, _best_diff = None, None
            for _i, _p in enumerate(_periods_raw):
                import re as _r3
                _m3 = _r3.search(r"(\d{4})", str(_p) or "")
                if not _m3: continue
                _yr = int(_m3.group(1))
                _ni_v = _ni_series[_i] if _i < len(_ni_series) else None
                if not isinstance(_ni_v, (int, float)) or _ni_v <= 0:
                    continue
                _d = abs(_yr - _cur_yr)
                if _best_diff is None or _d < _best_diff:
                    _best_idx, _best_diff = _i, _d
                    pe_fwd_year = _yr
            if _best_idx is not None:
                _ni_pick = float(_ni_series[_best_idx])
                _eps_synth = _ni_pick / _shares
                if _eps_synth > 0:
                    _pe_synth = float(last_price) / _eps_synth
                    if 1 < _pe_synth < 200:
                        pe_fwd = _pe_synth

    # Performance deltas — try Yahoo's hist_prices "perf_*" keys first;
    # fall back to MS price_performance block (perf_1d_pct etc.) when the
    # canonical historical_prices is empty (yfinance-blocked tickers).
    _ms_perf = (ms_price_performance or {}).get("performance") or {} \
        if isinstance(ms_price_performance, dict) else {}
    _ms_perf_keymap = {
        "perf_1d":  "perf_1d_pct",
        "perf_1w":  "perf_1w_pct",
        "perf_1m":  "perf_1m_pct",
        "perf_3m":  "perf_3m_pct",
        "perf_6m":  "perf_6m_pct",
        "perf_ytd": "perf_ytd_pct",
    }

    # Website-style perf + range from the LIVE price vs the close series —
    # this matches what the analyst sees on Investing/MS (their pre-computed
    # perf API lags the live quote).
    try:
        _live_price = float(last_price) if last_price is not None else None
    except (TypeError, ValueError):
        _live_price = None
    _series = hist_prices.get("close_series") if isinstance(hist_prices, dict) else None
    _live_perf = _perf_and_range_from_series(_series, _live_price)

    def _perf(key):
        # 1) website-style computation from the live price + series (matches the site)
        if isinstance(_live_perf.get(key), (int, float)):
            return float(_live_perf[key])
        # 2) the source's pre-computed perf (can lag the site)
        if isinstance(hist_prices, dict):
            v = hist_prices.get(key)
            if v is not None:
                try: return float(v)
                except (TypeError, ValueError): pass
        # 3) MarketScreener perf block
        ms_key = _ms_perf_keymap.get(key)
        if ms_key and isinstance(_ms_perf.get(ms_key), (int, float)):
            return float(_ms_perf[ms_key])
        return None

    # 52-week range — prefer the full-year series min/max computed above;
    # fall back to the source's explicit 52w fields only when the series
    # didn't span a year.
    low = high = current = None
    if _live_perf.get("range_52w_low") is not None:
        low, high = _live_perf["range_52w_low"], _live_perf["range_52w_high"]
    elif isinstance(hist_prices, dict):
        low = hist_prices.get("range_52w_low")
        high = hist_prices.get("range_52w_high")
    # Displayed price = the trust-laddered canonical price, EXCEPT when that
    # price comes from a laggy source (Investing's current_price can trail its
    # own dated close — BKMB canonical was 0.412 while the real latest close
    # was 0.400) and the series' dated last close disagrees by >1.5%. Then the
    # dated series close is the fresher truth. We never override a price that
    # came from Yahoo / yfinance-live (the freshest live source), so names where
    # Yahoo leads the Investing series (1180.SR) keep their live price.
    current = _live_price
    _cp_src = ((cv.get("current_price").canonical_source or "").lower()
               if cv.get("current_price") else "")
    _series_close = _live_perf.get("latest_close")
    if (isinstance(_series_close, (int, float)) and _series_close > 0
            and isinstance(_live_price, (int, float)) and _live_price > 0
            and "yahoo" not in _cp_src and "yfinance" not in _cp_src
            and abs(_series_close - _live_price) / _live_price > 0.015):
        current = _series_close
    if low is None or high is None or current is None:
        # Fallback: if we have only current, plot a degenerate range
        if current is not None:
            low = low if low is not None else current * 0.9
            high = high if high is not None else current * 1.1
        else:
            low, high, current = 0.0, 1.0, 0.5

    # MarketScreener's target_price comes as a dict
    # {"mean":29.86, "high":35.0, "low":26.5, "n_analysts":18}.
    # Yahoo / other providers may give a single float. Normalise.
    target_mean = None
    if isinstance(target, dict):
        target_mean = target.get("mean")
    elif isinstance(target, (int, float)):
        target_mean = target

    upside_pct = None
    try:
        if target_mean is not None and current is not None and current > 0:
            upside_pct = (float(target_mean) / current - 1.0) * 100
    except (TypeError, ValueError):
        pass

    target_fmt = "—"
    try:
        if target_mean is not None:
            _tv = float(target_mean)
            _tdp = _price_dp(_tv)
            target_fmt = (
                f"{currency} {_tv:,.{_tdp}f}" if currency
                else f"{_tv:,.{_tdp}f}"
            )
    except (TypeError, ValueError):
        pass

    # Compose meta-line pieces with safe fallbacks
    return SnapshotData(
        company_name=(profile.get("name") if isinstance(profile, dict) else None) or ticker,
        ticker=ticker,
        sector=(profile.get("sector") if isinstance(profile, dict) else None) or "—",
        industry=(profile.get("industry") if isinstance(profile, dict) else None) or "—",
        exchange=(profile.get("exchange") if isinstance(profile, dict) else None) or "—",
        period_label=period_label,
        rating=rating,
        n_analysts=n_analysts,
        target_price_fmt=target_fmt,
        upside_pct=upside_pct,
        last_close_fmt=_money(current),
        last_close_value=(float(current) if isinstance(current, (int, float)) else None),
        market_cap_fmt=_mc(mcap),
        report_date=_normalize_date(report_date),
        pe_fy_est_fmt=("—" if pe_fwd is None else f"{float(pe_fwd):.1f}x"),
        pe_fy_est_label=(f"P/E (FY{pe_fwd_year % 100:02d}E)"
                          if pe_fwd is not None and pe_fwd_year
                          else "P/E (FY EST)"),
        div_yield_fmt=_pct(div_yield),
        currency=currency or "",
        perf_1d=_perf("perf_1d"),
        perf_1w=_perf("perf_1w"),
        perf_1m=_perf("perf_1m"),
        perf_3m=_perf("perf_3m"),
        perf_6m=_perf("perf_6m"),
        perf_ytd=_perf("perf_ytd"),
        range_low=float(low), range_high=float(high), range_current=float(current),
        highlights=highlights or _llm_highlights_or_template(
            ticker=ticker,
            cv=cv, currency=currency, current_price=current,
            mcap=mcap, target_mean=target_mean, upside_pct=upside_pct,
            pe_fwd=pe_fwd, div_yield=div_yield,
            range_low=low, range_high=high,
            n_analysts=n_analysts, rs=rating_split,
        ),
        sources_line=_sources_line(cv),
        analyst_name=analyst_name,
        gen_date=datetime.utcnow().strftime("%d %b %Y"),
        # Price as-of: prefer the Investing historical series (same data
        # vintage as the displayed Investing equity price) over whatever is
        # in canonical historical_prices, which may be an iShares EEM proxy
        # that is fresh-dated but UNRELATED to this stock's price age.
        freshness_banner=_compute_freshness_banner(
            cv, live_quote_record,
            price_asof=(_price_asof_from_history(historical_override)
                        or _price_asof_from_history(hist_prices))),
    )
