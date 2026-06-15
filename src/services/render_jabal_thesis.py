"""
Jabal — Slide 2 (Thesis & Expectations) renderer.

Layout (per design_spec.md):
  1. Header strip
  2. Section label (INVESTMENT THESIS) + Georgia 17 title ("Executive Summary")
  3. Body card — 4-6 sentence thesis paragraph
  4. Section label (Q2 2026 EARNINGS EXPECTATIONS) + table:
        Metric | Jabal Est. | Consensus | Δ | YoY%
        rows for Revenue / EBITDA / EBITDA margin / Net income / EPS / Dividend
  5. Catalysts + Key Risks — two side-by-side cards (3 bullets each)
  6. Numbered "What to watch on the print"
  7. Footer

Data sources used:
  - canonical_store.valuation_forward            → consensus EPS/PE/etc.
  - canonical_store.valuation_historical         → YoY base
  - get_observations_by_provider(commodities)    → commodity context for thesis text
  - get_observations_by_provider(macro)          → IMF/WB context for thesis text
  - canonical_store.rating_split                 → consensus colour in opening sentence
"""

from __future__ import annotations

from dataclasses import dataclass, field
from typing import Optional

from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from src.services.jabal_design_tokens import (
    BLACK, GRAY, MUTED, GOLD, POS, NEG, CARD, WHITE,
    FONT_DISPLAY, FONT_UI,
    SZ_SECTION, SZ_KICKER, SZ_VALUE, SZ_VALUE_LG, SZ_LABEL, SZ_BODY,
    SZ_META, SZ_HEADER, SZ_FOOTER, SZ_BULLET_PILL, SZ_TAB_NUM, SZ_TINY,
    PAGE_W_IN, PAGE_H_IN, MARGIN_L, MARGIN_R, CONTENT_W,
    RULE_THICK_PT, BORDER_THICK_PT, LEFT_ACCENT_W_IN,
    in_, signed_color,
)
from src.services.canonical_store import (
    get_all_fields, get_observations_by_provider,
)
from src.services.render_jabal_snapshot import (
    _text, _hrule, _card, _section_label, _header_strip, _footer,
)


# ── Slide 2 sections ──────────────────────────────────────────

def _section_hero(slide, top: float, label: str, title: str):
    """Section label + Georgia 17 title underneath, used for slide 2/3 heroes."""
    _hrule(slide, MARGIN_L, top, CONTENT_W, color=MUTED)
    _text(slide, MARGIN_L, top + 0.10, CONTENT_W, 0.22, label,
          size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)
    _text(slide, MARGIN_L, top + 0.32, CONTENT_W, 0.50, title,
          font=FONT_DISPLAY, size=SZ_SECTION, color=BLACK)


def _body_card(slide, left: float, top: float, width: float, height: float,
                body: str):
    """Cream-fill card with gold left-accent and the thesis paragraph inside."""
    _card(slide, left, top, width, height, fill=CARD, border=MUTED,
           left_accent=GOLD)
    # Inset 0.20 from accent edge
    _text(slide, left + 0.20, top + 0.10, width - 0.32, height - 0.20,
          body, size=SZ_BODY, color=BLACK)


def _estimates_table(slide, top: float, rows: list[dict],
                      period_label: str = "ESTIMATE"):
    """5-column estimates table: METRIC | Jabal | YoY | QoQ | CONSENSUS.

    `period_label` is the dynamic header for the Jabal-estimate column
    (e.g. 'Q2 2026E'). Column order matches the institutional reference
    layout: analyst estimate first, deltas in the middle, consensus last.

    Row dict shape: {metric, jabal, yoy, qoq, consensus, is_margin}.
    `is_margin` renders YoY/QoQ in basis points; other rows render in %.
    """
    headers = ["METRIC", period_label.upper(), "YoY", "QoQ", "CONSENSUS"]
    col_w   = [2.20, 1.10, 1.00, 1.00, 1.30]
    row_h   = 0.30
    header_top = top
    # Header
    x = MARGIN_L
    for i, h in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        _text(slide, x, header_top, col_w[i] - 0.05, row_h, h,
              size=SZ_LABEL, color=MUTED, all_caps=True, align=align, wrap=False)
        x += col_w[i]
    _hrule(slide, MARGIN_L, header_top + row_h - 0.02, CONTENT_W,
            color=MUTED)
    # Body rows
    for ri, row in enumerate(rows):
        y = header_top + row_h + ri * row_h
        if ri % 2 == 1:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                in_(MARGIN_L), in_(y - 0.02),
                in_(CONTENT_W), in_(row_h))
            band.fill.solid(); band.fill.fore_color.rgb = CARD
            band.line.fill.background()
        is_margin = bool(row.get("is_margin"))
        x = MARGIN_L
        for i, key in enumerate(["metric", "jabal", "yoy", "qoq", "consensus"]):
            val = row.get(key, "—")
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            color = BLACK
            if key in ("yoy", "qoq") and isinstance(val, (int, float)):
                color = signed_color(val)
                # Margin row: render in basis points. Other rows: percent.
                if is_margin:
                    val = f"{val * 100:+.0f} bps"   # input value is in pp; ×100 → bps
                elif abs(val) > 300:
                    # Off a near-zero base (e.g. loss → profit) the % is real
                    # but not meaningful; institutional notation is "n/m".
                    val = "n/m"; color = BLACK
                else:
                    val = f"{val:+.1f}%"
            elif val is None:
                val = "—"
            _text(slide, x, y, col_w[i] - 0.05, row_h, str(val),
                  size=SZ_BODY, color=color, align=align, wrap=False)
            x += col_w[i]


def _annual_estimates_table(slide, top: float, rows: list[dict],
                              fy_labels: list[str]):
    """Annual-FY fallback table — columns: METRIC | FY+1 | FY+2 | FY+3 |
    YoY (FY+1 vs last actual) | CAGR (FY+3 vs last actual, annualised).

    Row dict shape: {metric, fy1, fy2, fy3, yoy, cagr}. YoY + CAGR
    rendered with signed colour; currency cells rendered as black text.

    Mohamed (2026-05): "instead of the last column, or in addition we can
    add the cagr since it shows a forward of three years." Kept both —
    YoY anchors against the prior FY actual (short-term momentum read),
    CAGR captures the multi-year trajectory implied by the strip.
    """
    fy = list(fy_labels) + ["—"] * (3 - len(fy_labels))
    headers = ["METRIC", fy[0].upper(), fy[1].upper(), fy[2].upper(),
                f"YoY {fy[0]}", f"3Y CAGR"]
    # Columns MUST sum to CONTENT_W (6.60") or the 6th column (3Y CAGR)
    # runs off the right edge. 1.80 + 0.88·3 + 1.08·2 = 6.60.
    col_w   = [1.80, 0.88, 0.88, 0.88, 1.08, 1.08]
    row_h   = 0.30
    header_top = top
    x = MARGIN_L
    for i, h in enumerate(headers):
        align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
        _text(slide, x, header_top, col_w[i] - 0.05, row_h, h,
              size=SZ_LABEL, color=MUTED, all_caps=True, align=align, wrap=False)
        x += col_w[i]
    _hrule(slide, MARGIN_L, header_top + row_h - 0.02, CONTENT_W, color=MUTED)
    for ri, row in enumerate(rows):
        y = header_top + row_h + ri * row_h
        if ri % 2 == 1:
            band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                in_(MARGIN_L), in_(y - 0.02),
                in_(CONTENT_W), in_(row_h))
            band.fill.solid(); band.fill.fore_color.rgb = CARD
            band.line.fill.background()
        x = MARGIN_L
        for i, key in enumerate(["metric", "fy1", "fy2", "fy3", "yoy", "cagr"]):
            val = row.get(key, "—")
            align = PP_ALIGN.LEFT if i == 0 else PP_ALIGN.RIGHT
            color = BLACK
            if key in ("yoy", "cagr") and isinstance(val, (int, float)):
                color = signed_color(val)
                val = "n/m" if abs(val) > 300 else f"{val:+.1f}%"
                if val == "n/m":
                    color = BLACK
            elif val is None:
                val = "—"
            _text(slide, x, y, col_w[i] - 0.05, row_h, str(val),
                  size=SZ_BODY, color=color, align=align, wrap=False)
            x += col_w[i]


def _two_col_pillared_card(slide, top: float, height: float,
                              left_title: str, left_bullets: list[str],
                              right_title: str, right_bullets: list[str]):
    """Side-by-side cards (Catalysts | Key Risks)."""
    card_w = (CONTENT_W - 0.20) / 2
    # LEFT
    _card(slide, MARGIN_L, top, card_w, height, fill=WHITE, border=MUTED,
           left_accent=POS)
    _text(slide, MARGIN_L + 0.18, top + 0.08, card_w - 0.20, 0.30,
          left_title, size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)
    bullet_y = top + 0.50
    for b in left_bullets[:3]:
        _bullet_dot(slide, MARGIN_L + 0.20, bullet_y + 0.02, color=POS)
        _text(slide, MARGIN_L + 0.42, bullet_y - 0.02, card_w - 0.50, 0.42, b,
              size=SZ_BODY, color=BLACK)
        bullet_y += 0.46
    # RIGHT
    r_left = MARGIN_L + card_w + 0.20
    _card(slide, r_left, top, card_w, height, fill=WHITE, border=MUTED,
           left_accent=NEG)
    _text(slide, r_left + 0.18, top + 0.08, card_w - 0.20, 0.30,
          right_title, size=SZ_KICKER, color=GRAY, all_caps=True, bold=True)
    bullet_y = top + 0.50
    for b in right_bullets[:3]:
        _bullet_dot(slide, r_left + 0.20, bullet_y + 0.02, color=NEG)
        _text(slide, r_left + 0.42, bullet_y - 0.02, card_w - 0.50, 0.42, b,
              size=SZ_BODY, color=BLACK)
        bullet_y += 0.46


def _bullet_dot(slide, left, top, *, color=GOLD, size=0.12):
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                  in_(left), in_(top),
                                  in_(size), in_(size))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    return dot


def _numbered_list(slide, top: float, items: list[str]):
    """1/2/3 numerals in gold + body text. Used for 'What to watch on the print'."""
    row_h = 0.27
    for i, body in enumerate(items[:5]):
        y = top + i * row_h
        _text(slide, MARGIN_L, y, 0.25, 0.22, str(i + 1),
              size=SZ_BODY, color=GOLD, bold=True)
        _text(slide, MARGIN_L + 0.30, y, CONTENT_W - 0.30, 0.22, body,
              size=SZ_BODY, color=BLACK)


# ── Public entry point ────────────────────────────────────────

@dataclass
class ThesisData:
    """Slide 2 inputs. Built by the orchestrator from canonical_store +
    light templating of commodities / macro context."""
    exec_summary_body: str
    estimates_rows: list[dict]
    estimates_footnote: str
    estimates_subtitle: str = ""  # Subtitle above the table, e.g. "Jabal estimates vs. consensus  ·  SAR millions unless stated"
    estimates_period_label: str = "ESTIMATE"  # Column header for the Jabal-estimate column, e.g. "Q2 2026E"
    catalysts: list[str] = field(default_factory=list)
    risks: list[str] = field(default_factory=list)
    watch_list: list[str] = field(default_factory=list)
    sources_line: str = ""
    analyst_name: str = "Jabal Research"
    gen_date: str = ""
    total_pages: int = 3
    period_heading: str = "Earnings Expectations"
    # Annual-FY fallback fields — populated when per-quarter consensus is
    # only available as the MS-annual÷4 synthetic proxy and the analyst
    # call is to suppress the unreliable quarterly view.
    is_annual_table: bool = False
    annual_rows: list[dict] = field(default_factory=list)
    annual_fy_labels: list[str] = field(default_factory=list)


def render_thesis_slide(prs, data: ThesisData):
    blank = next((L for L in prs.slide_layouts if L.name.lower() == "blank"),
                  prs.slide_layouts[-1])
    slide = prs.slides.add_slide(blank)

    _header_strip(slide, 2, "Thesis & Expectations")

    # Investment Thesis hero + body card
    _section_hero(slide, 1.08, "Investment Thesis", "Executive Summary")
    _body_card(slide, MARGIN_L, 1.96, CONTENT_W, 1.85,
                data.exec_summary_body)

    # Period-aware section label. Falls back to a generic title when the
    # orchestrator didn't supply a quarter (e.g. carry-forward case).
    _section_label(slide, MARGIN_L, 3.96, CONTENT_W, data.period_heading or "Earnings Expectations")
    _text(slide, MARGIN_L, 4.26, CONTENT_W, 0.18,
          data.estimates_subtitle or "Jabal estimates vs. consensus  ·  Local currency unless stated",
          size=Pt(9), color=GRAY)
    if data.is_annual_table and data.annual_rows:
        _annual_estimates_table(slide, 4.48, data.annual_rows,
                                  fy_labels=data.annual_fy_labels or [])
    else:
        _estimates_table(slide, 4.48, data.estimates_rows,
                           period_label=data.estimates_period_label or "ESTIMATE")
    _text(slide, MARGIN_L, 6.88, CONTENT_W, 0.18,
          data.estimates_footnote,
          size=Pt(9), color=GRAY)

    # Catalysts + Risks
    _two_col_pillared_card(
        slide, 7.38, 1.95,
        "Catalysts", data.catalysts,
        "Key Risks", data.risks,
    )

    # Numbered list
    _section_label(slide, MARGIN_L, 9.53, CONTENT_W, "What to Watch on the Print")
    _numbered_list(slide, 9.81, data.watch_list)

    _footer(slide, 2, data.total_pages, data.sources_line,
             data.analyst_name, data.gen_date)
    return slide


# ── Data adapter ──────────────────────────────────────────────

# Per-sector operating levers for the LLM-absent fallback. The OPENING names
# these (instead of the generic "sector-specific operating levers") so the
# scaffold reads like a specialist wrote it even when Gemini is unavailable.
_SECTOR_LEVERS = {
    "bank": "net interest income, loan growth, fee income, and credit quality",
    "energy": "production volumes, realized prices, and lifting costs",
    "materials": "product spreads, sales volumes, and feedstock costs",
    "tech": "revenue mix, margin trajectory, and forward guidance",
    "telco": "service revenue, subscriber/ARPU trends, and capex intensity",
    "industrial": "order intake, margin trajectory, and input costs",
    "consumer": "volume/price mix, gross margin, and demand trends",
    "healthcare": "revenue growth, margin trajectory, and pipeline progress",
    "utility": "tariff/volume trends, fuel costs, and capex execution",
}


def _sector_key(sector_l: str, industry_l: str) -> str:
    s = f"{sector_l} {industry_l}"
    if "bank" in s or "financ" in s or "insur" in s:
        return "bank"
    if "oil" in s or "gas" in s or "energy" in s:
        return "energy"
    if "chemical" in s or "materi" in s or "mining" in s or "metal" in s:
        return "materials"
    if "internet" in s or "software" in s or "semic" in s or "technolog" in s:
        return "tech"
    if "telecom" in s or "communicat" in s:
        return "telco"
    if "industrial" in s or "machin" in s or "aerospace" in s:
        return "industrial"
    if "consumer" in s or "retail" in s or "food" in s or "bever" in s:
        return "consumer"
    if "health" in s or "pharma" in s or "medic" in s:
        return "healthcare"
    if "utilit" in s or "power" in s or "electric" in s:
        return "utility"
    return ""


def _template_exec_summary(cv: dict, commodities: dict,
                              macro_obs: dict, ticker: str = "") -> str:
    """Compose a 4-sentence thesis paragraph from canonical data + grounded
    FY actuals. This is the deterministic fallback used when the LLM is
    unavailable (e.g. Gemini billing/quota); it names the sector's operating
    levers and cites a verified full-year actual so it reads like a specialist
    wrote it, not a blank template."""
    name = "the company"
    sector = industry = "—"
    profile = cv.get("company_profile")
    if profile and isinstance(profile.value, dict):
        name = profile.value.get("name") or name
        sector = profile.value.get("sector") or sector
        industry = profile.value.get("industry") or industry

    rating_val = cv.get("rating_split")
    rating = ""
    n_an = 0
    if rating_val and isinstance(rating_val.value, dict):
        # Prettify the provider enum so the fallback paragraph reads
        # "Strong Buy" rather than "strong_buy" / "STRONG_BUY".
        from src.services.render_jabal_snapshot import _pretty_rating
        rating = _pretty_rating(rating_val.value.get("consensus")) or ""
        n_an   = int(rating_val.value.get("total") or 0)

    target = cv.get("target_price")
    target_text = ""
    if target and isinstance(target.value, dict):
        m = target.value.get("mean")
        if m:
            target_text = f"; consensus target {m:.2f}"

    val_hist = cv.get("valuation_historical")
    pe_text = ""
    if val_hist and isinstance(val_hist.value, dict):
        pe = val_hist.value.get("pe")
        if isinstance(pe, list) and any(pe):
            recent = [p for p in pe if isinstance(p, (int, float))]
            if recent:
                pe_text = f" The shares trade at a P/E around {recent[-1]:.1f}x trailing earnings."

    commodity_text = ""
    industry_commodities = (commodities.get("company_profile") or {}).get("industry_commodities", {}) \
        if commodities else {}
    if industry_commodities:
        bits = []
        for tag, info in industry_commodities.items():
            if not isinstance(info, dict):
                continue
            val = info.get("value")
            yoy = info.get("yoy_pct")
            unit = info.get("unit", "")
            if val is None:
                continue
            yoy_str = f" ({yoy:+.1f}% YoY)" if yoy is not None else ""
            bits.append(f"{tag} at {val:.0f} {unit}{yoy_str}")
        if bits:
            commodity_text = " The macro and commodity backdrop is anchored by " \
                + "; ".join(bits[:2]) + "."

    # Macro template: prefer IMF current-year + next-year forecasts over
    # World Bank historical actuals. WB falls back only when IMF lacks a
    # series for the country. Every figure year-stamped so the resulting
    # sentence reads "IMF 2026 forecast 3.5%" rather than the misleading
    # "GDP recently at 1.6%" (which was a stale 2024 actual).
    macro_text = ""
    mp = macro_obs.get("company_profile") if macro_obs else None
    if mp:
        gdp_now      = mp.get("gdp_growth_fcst_pct") or mp.get("gdp_growth_pct")
        gdp_now_yr   = mp.get("gdp_growth_fcst_year") or mp.get("macro_year") or ""
        gdp_now_src  = "IMF" if mp.get("gdp_growth_fcst_pct") is not None else "WB"
        gdp_next     = mp.get("gdp_growth_fcst_next_pct")
        gdp_next_yr  = mp.get("gdp_growth_fcst_next_year") or ""
        infl_now     = mp.get("inflation_fcst_pct") or mp.get("inflation_pct")
        infl_now_yr  = mp.get("inflation_fcst_year") or mp.get("macro_year") or ""
        infl_now_src = "IMF" if mp.get("inflation_fcst_pct") is not None else "WB"
        parts = []
        if gdp_now is not None:
            parts.append(f"GDP growth {gdp_now:.1f}% ({gdp_now_src} {gdp_now_yr})")
        if gdp_next is not None and gdp_next_yr:
            parts.append(f"IMF {gdp_next_yr} forecast {gdp_next:.1f}%")
        if infl_now is not None:
            parts.append(f"inflation {infl_now:.1f}% ({infl_now_src} {infl_now_yr})")
        if parts:
            macro_text = " Macro context: " + ", ".join(parts) + "."

    rating_line = ""
    if rating and n_an:
        rating_line = f" Street consensus is {rating} ({n_an} analysts covering){target_text}."

    # Sector-aware closing — feedstock/H2 demand only makes sense for
    # industrial / chemical names; banks, internet, etc. need their own
    # framing. Keeps the LLM-absent fallback honest rather than
    # generic-industrial regardless of sector.
    sector_l = (sector or "").lower()
    industry_l = (industry or "").lower()
    if "bank" in sector_l or "bank" in industry_l or "financial" in sector_l:
        closing = (
            " The print hinges on net interest income, loan growth, "
            "credit quality, and capital deployment; the read on forward "
            "guidance is the swing factor."
        )
    elif "oil" in industry_l or "gas" in industry_l or "energy" in sector_l:
        closing = (
            " The print hinges on production volumes, realized prices, "
            "and lifting costs; management's tone on capex and project "
            "ramps is the swing factor."
        )
    elif "internet" in industry_l or "software" in industry_l or "technology" in sector_l:
        closing = (
            " The print hinges on revenue mix, margin trajectory, and "
            "guidance; product / regulatory updates are the swing factor."
        )
    elif "mining" in industry_l or "metals" in industry_l:
        closing = (
            " The print hinges on production volumes, realized commodity "
            "prices, and unit costs; guidance and project execution are "
            "the swing factor."
        )
    else:
        closing = (
            " The print hinges on revenue growth, margin trajectory, and "
            "guidance; management's tone on forward demand is the swing factor."
        )

    # Name the sector's operating levers in the opening (no "sector-specific
    # operating levers" placeholder).
    sk = _sector_key(sector_l, industry_l)
    levers = _SECTOR_LEVERS.get(sk, "revenue growth, margin trajectory, and guidance")
    opening = (f"{name} enters the upcoming print with the Street focused on "
               f"{levers}.")

    # Cite a VERIFIED full-year actual from the grounding store so the
    # fallback is grounded, not generic ("supported by FY2025's +13.3% net
    # profit growth, with a 13.57% ROE").
    grounded_text = ""
    try:
        from src.services.disclosed_loader import load_disclosed
        fy = (load_disclosed(ticker) or {}).get("fy_highlights") or {} if ticker else {}
        if fy:
            period = fy.get("period") or "the prior year"
            npg = fy.get("net_profit_growth_pct")
            roe = fy.get("roe_pct")
            rev_g = fy.get("revenue_growth_pct") or fy.get("total_income_growth_pct")
            bits = []
            if isinstance(npg, (int, float)):
                bits.append(f"{period}'s {npg:+.1f}% net-profit growth")
            elif isinstance(rev_g, (int, float)):
                bits.append(f"{period}'s {rev_g:+.1f}% revenue growth")
            if isinstance(roe, (int, float)):
                bits.append(f"a {roe:.1f}% ROE")
            if bits:
                grounded_text = (" Recent delivery is anchored by "
                                 + " and ".join(bits) + ".")
    except Exception:
        grounded_text = ""

    body = (
        opening
        + grounded_text
        + rating_line
        + pe_text
        + commodity_text
        + macro_text
        + closing
    )
    return body


def _fmt_money_b(v):
    """Render a raw money value as 'X.YB' / 'X.YT' / 'X,Y00M'. None on bad input."""
    if not isinstance(v, (int, float)):
        return None
    if abs(v) >= 1e12: return f"{v/1e12:,.2f}T"
    if abs(v) >= 1e9:  return f"{v/1e9:,.1f}B"
    if abs(v) >= 1e6:  return f"{v/1e6:,.0f}M"
    return f"{v:,.0f}"


def _prior_year_same_q(period_label: str) -> str | None:
    """'2024-Q3' -> '2023-Q3'; '2024-Q3 (Mar)' -> '2023-Q3 (Mar)'. None on bad input."""
    import re as _re
    m = _re.match(r"^\s*(\d{4})(\D.*)$", str(period_label or ""))
    if not m:
        return None
    try:
        return f"{int(m.group(1)) - 1}{m.group(2)}"
    except ValueError:
        return None


def _yoy_pct(curr, prev) -> float | None:
    try:
        if curr is None or prev in (None, 0):
            return None
        return (float(curr) / float(prev) - 1.0) * 100
    except (TypeError, ValueError, ZeroDivisionError):
        return None


def _yoy_bps(curr_pct, prev_pct) -> float | None:
    """Margin YoY in basis points (return as percentage for unified rendering)."""
    try:
        if curr_pct is None or prev_pct is None:
            return None
        return float(curr_pct) - float(prev_pct)
    except (TypeError, ValueError):
        return None


def _ms_quarterly_split(ms_q: dict | None) -> tuple[dict, dict, dict, dict, dict]:
    """From `ms_quarterly_forecasts.quarterly`, return (next_est,
    latest_actual, prior_year_actual, prior_of_next, prior_quarter).
    Each is a dict {metric_key: value}.

    MS interleaves actuals and forecasts in one period list. We split by
    whether the announcement_date is in the past (actual) or future
    (estimate). Falls back to empty dicts on any shape problem.
    """
    empty = ({}, {}, {}, {}, {})
    if not isinstance(ms_q, dict):
        return empty
    q = ms_q.get("quarterly") or {}
    if not isinstance(q, dict):
        return empty
    periods   = q.get("periods")  or []
    net_sales = q.get("net_sales") or []
    ebitda    = q.get("ebitda")   or []
    nii       = q.get("nii")      or []   # banks (when MS publishes it)
    net_inc   = q.get("net_income") or []
    eps       = q.get("eps")      or []
    ann       = q.get("announcement_dates") or []
    n = len(periods)
    if not n:
        return empty
    # Pad short lists with None so zip aligns.
    def _pad(xs): return list(xs) + [None] * (n - len(xs))
    rows = list(zip(periods, _pad(net_sales), _pad(ebitda), _pad(nii),
                     _pad(net_inc), _pad(eps), _pad(ann)))

    from datetime import datetime as _dt
    today = _dt.now().date()

    def _parse_ms_date(date_str: str):
        """MS publishes announcement dates as MM/DD/YY (e.g. '4/23/26').
        Try each known shape in turn and return a `date` or None."""
        s = str(date_str or "").strip()
        if not s or s == "-":
            return None
        for fmt in ("%Y-%m-%d", "%m/%d/%y", "%m/%d/%Y", "%-m/%-d/%y"):
            try:
                return _dt.strptime(s, fmt).date()
            except (ValueError, TypeError):
                continue
        return None

    def _is_estimate(date_str) -> bool:
        """True if the row is a forward forecast (no past announcement)."""
        d = _parse_ms_date(date_str)
        if d is None:
            # Empty / '-' / unparseable → treat as forward estimate.
            return True
        return d >= today

    actuals = [r for r in rows if not _is_estimate(r[6])]
    estimates = [r for r in rows if _is_estimate(r[6])]

    def _to_dict(r):
        p, rev, eb, n_, ni, e, ad = r
        return {
            "period": p, "revenue": rev, "ebitda": eb,
            "nii": n_, "net_income": ni, "eps": e, "ann": ad,
        }

    next_est = _to_dict(estimates[0]) if estimates else {}
    latest_actual = _to_dict(actuals[-1]) if actuals else {}

    # Helper: parse a period label ("Q2 2026" / "2Q26" / "2026-Q2" / "2026Q2")
    # into (year, quarter). Returns None on no match.
    import re as _r
    def _parse_q(period: str):
        s = str(period or "")
        m = (_r.search(r"(\d{4})\s*Q(\d)|Q(\d)\s*(\d{4})", s, _r.I)
             or _r.search(r"(\d{4})-Q(\d)", s, _r.I)
             or _r.search(r"(\d{4})Q(\d)", s, _r.I))
        if not m:
            return None
        y = int(m.group(1) or m.group(4))
        q = int(m.group(2) or m.group(3))
        return (y, q)

    def _find_in_actuals(target_year, target_q):
        for a in actuals:
            yq = _parse_q(a[0])
            if yq == (target_year, target_q):
                return _to_dict(a)
        return None

    # Prior-year-same-Q for the LATEST actual (used for last-reported YoY).
    prior = {}
    if latest_actual.get("period"):
        yq = _parse_q(latest_actual["period"])
        if yq:
            prior = _find_in_actuals(yq[0] - 1, yq[1]) or {}

    # Prior-year-same-Q for the NEXT estimate (used for forecast-YoY).
    prior_of_next = {}
    if next_est.get("period"):
        yq = _parse_q(next_est["period"])
        if yq:
            prior_of_next = _find_in_actuals(yq[0] - 1, yq[1]) or {}

    # Prior-quarter actual: the quarter immediately preceding the next
    # forecast. Needed for QoQ (e.g. Q2 2026 forecast vs Q1 2026 actual).
    prior_quarter = {}
    if next_est.get("period"):
        yq = _parse_q(next_est["period"])
        if yq:
            ny, nq = yq
            if nq > 1:
                prior_quarter = _find_in_actuals(ny, nq - 1) or {}
            else:
                prior_quarter = _find_in_actuals(ny - 1, 4) or {}

    return next_est, latest_actual, prior, prior_of_next, prior_quarter


def _investing_actuals_yoy(ticker: str) -> dict:
    """Pull NEXT forecast vs prior-year-same-quarter actual from Investing's
    earnings page, returning {revenue, eps} YoY percentages.

    Third-tier fallback for the slide-2 forecast-YoY column when MS doesn't
    publish a clean quarterly forecast block. Forecast row is the first
    earnings_history entry where epsActual is None (the upcoming print);
    prior_actual is the entry with reportYear = forecast_year - 1 and the
    same quarter.

    Falls back to (latest actual vs prior-year actual) only when no
    forecast row exists at all (rare; some thinly-covered names have only
    actuals on the Investing page).
    """
    try:
        from src.providers.probe_investing import _fetch_earnings_page, _slug
    except ImportError:
        return {}
    slug = _slug(ticker)
    if not slug:
        return {}
    state = _fetch_earnings_page(slug)
    if not state:
        return {}
    es = state.get("earningsStore") or {}
    if not isinstance(es, dict):
        return {}
    rows = es.get("earnings") or []
    actuals = [r for r in rows
                if isinstance(r, dict)
                and isinstance(r.get("epsActual"), (int, float))
                and isinstance(r.get("reportYear"), int)
                and isinstance(r.get("reportMonth"), int)]
    # Forecast = first row with epsActual None / epsForecast present.
    forecast = None
    for r in rows or []:
        if not isinstance(r, dict): continue
        if r.get("epsActual") is None and isinstance(r.get("reportYear"), int):
            if (isinstance(r.get("epsForecast"), (int, float))
                or isinstance(r.get("revenueForecast"), (int, float))):
                forecast = r
                break

    # Prefer forecast-vs-prior-year-actual when both exist. Falls back to
    # latest-actual-vs-prior-year-actual when no forecast on the page.
    if forecast:
        ny = forecast["reportYear"]
        nq = ((forecast["reportMonth"] - 1) // 3 + 1)
        prior = next((
            r for r in actuals
            if r["reportYear"] == ny - 1
            and ((r["reportMonth"] - 1) // 3 + 1) == nq
        ), None)
        if not prior:
            return {}
        latest = forecast
        # Forecast row uses different key names — adapt below.
        rev_l = forecast.get("revenueForecast")
        rev_p = prior.get("revenueActual")
        eps_l = forecast.get("epsForecast")
        eps_p = prior.get("epsActual")
        out = {}
        if isinstance(rev_l, (int, float)) and isinstance(rev_p, (int, float)) and rev_p:
            out["revenue"] = (rev_l / rev_p - 1.0) * 100
        if isinstance(eps_l, (int, float)) and isinstance(eps_p, (int, float)) and eps_p:
            out["eps"] = (eps_l / eps_p - 1.0) * 100
        return out

    # No forecast row — fall back to last-actual-vs-prior-year-actual.
    if not actuals:
        return {}
    latest = actuals[0]
    latest_year = latest["reportYear"]
    latest_qm = ((latest["reportMonth"] - 1) // 3 + 1)
    prior = next((
        r for r in actuals[1:]
        if r["reportYear"] == latest_year - 1
        and ((r["reportMonth"] - 1) // 3 + 1) == latest_qm
    ), None)
    if not prior:
        return {}
    out = {}
    rev_l, rev_p = latest.get("revenueActual"), prior.get("revenueActual")
    if isinstance(rev_l, (int, float)) and isinstance(rev_p, (int, float)) and rev_p:
        out["revenue"] = (rev_l / rev_p - 1.0) * 100
    eps_l, eps_p = latest.get("epsActual"), prior.get("epsActual")
    if isinstance(eps_l, (int, float)) and isinstance(eps_p, (int, float)) and eps_p:
        eps_yoy = (eps_l / eps_p - 1.0) * 100
        # Investing rounds bank EPS to 2dp (BKMB Q1 2026 = 0.01, Q1 2025 = 0.01
        # → YoY 0.0% even when Net Income grew 9%). When the rounded values
        # match exactly but the underlying business changed, fall back to
        # the Net-Income YoY as the EPS proxy — true absent share-count change.
        if eps_l == eps_p and abs(eps_l) <= 0.05:
            ni_yoy = None
            try:
                from src.services.store_actuals import latest_actuals  # unused but kept for safety
            except ImportError:
                pass
            # Compute NI YoY from the same Investing earnings page if it
            # exposes revenue (a reasonable proxy for NI YoY when bank EPS
            # is rounded). We don't have NI in earnings rows on Investing
            # — instead, just suppress the EPS YoY rather than report a
            # misleading 0.0%.
            out["eps"] = None  # caller renders '—'
        else:
            out["eps"] = eps_yoy
    return out


def _build_estimates_rows(cv: dict, quarterly: list | None = None,
                            is_bank: bool = False,
                            ms_quarterly_forecasts: dict | None = None,
                            ticker: str = "",
                            currency: str = "",
                            memo_data: dict | None = None) -> tuple[list[dict], str]:
    """Build rows for the slide-2 estimates table.

    Returns (rows, unit_suffix). The caller bakes unit_suffix into the
    table subtitle (e.g. "SAR millions unless stated").

    Row schema:
      Non-bank: Revenue / EBITDA / Net Income / EPS / EBITDA Margin
      Bank:     Operating Income / Net Income / EPS

    Columns:
      • JABAL EST   — '—' in the auto-deck; analyst fills in PPT.
      • YoY         — next-Q consensus vs prior-year-same-Q actual (%).
      • QoQ         — next-Q consensus vs immediately-prior-Q actual (%).
                       Margin rows render YoY/QoQ in bps via `is_margin`.
      • CONSENSUS   — next-Q consensus from Investing/MS.
    """
    val_fwd = cv.get("valuation_forward")
    fwd = val_fwd.value if val_fwd and isinstance(val_fwd.value, dict) else {}

    # Defer canonical_store reads until AFTER we've consulted MS quarterly.
    # MarketScreener publishes the entire quarterly forecast table (rev,
    # EBITDA, NI, EPS, all four metrics, with the same labels and units)
    # — when present, that table IS the consensus, full stop. Reading
    # `fwd.revenue_next_q` first lets Investing's value override MS, which
    # is the bug that produced SABIC's 3.01 SARB revenue (Investing) /
    # 0.94 SARB EBITDA (MS) — mixed sources giving a fake 31% margin and
    # garbage YoY/QoQ deltas.
    rev_q_consensus = None
    eps_q_consensus = None
    ebitda_q_consensus = None
    ni_q_consensus = None
    nii_q_consensus = None

    # MS /finances/ quarterly forecasts — the canonical_store doesn't carry
    # Revenue/EBITDA/NI consensus for the next quarter, but the MS payload
    # does. When the upstream pipeline passed it through, use it to populate
    # the CONSENSUS column AND to compute YoY (latest actual vs prior-year
    # same quarter from the same MS table). Falls back to canonical_store +
    # payload.quarterly_actuals otherwise.
    yoy_rev = yoy_ebitda = yoy_nii = yoy_ni = yoy_eps = yoy_margin = None
    qoq_rev = qoq_ebitda = qoq_nii = qoq_ni = qoq_eps = qoq_margin = None
    used_ms = False
    if ms_quarterly_forecasts:
        next_est, latest, prior, prior_of_next, prior_quarter = _ms_quarterly_split(ms_quarterly_forecasts)
        # MS publishes unit_scale as a *string* ("million", "billion",
        # "thousand"). Map it to a numeric multiplier so the formatter
        # renders absolute values (e.g. 2614M -> "2.6B").
        _UNIT_MULT = {
            "thousand": 1e3, "thousands": 1e3,
            "million":  1e6, "millions":  1e6, "m": 1e6,
            "billion":  1e9, "billions":  1e9, "b": 1e9,
        }
        raw_scale = ms_quarterly_forecasts.get("unit_scale")
        if isinstance(raw_scale, (int, float)):
            unit_scale = float(raw_scale) or 1.0
        else:
            unit_scale = _UNIT_MULT.get(str(raw_scale or "").strip().lower(), 1.0)
        def _scale(v):
            return v * unit_scale if isinstance(v, (int, float)) else v
        if next_est:
            # MS quarterly forecast WINS — when MS publishes a forward-
            # dated Q+1 row, every metric in that row is the canonical
            # consensus. Don't `or` against canonical_store (Investing-
            # sourced) values; mixing sources gave us SABIC's 3.01B
            # revenue (Investing) + 0.94B EBITDA (MS) → fake 31% margin.
            rev_q_consensus    = _scale(next_est.get("revenue"))
            ebitda_q_consensus = _scale(next_est.get("ebitda"))
            ni_q_consensus     = _scale(next_est.get("net_income"))
            nii_q_consensus    = _scale(next_est.get("nii"))
            eps_q_consensus    = next_est.get("eps")
        # ANALYTICAL CONTRACT: YoY in a "Q<n> Earnings Expectations" table
        # is the forecast-vs-prior-year-actual comparison — what the next
        # quarter's consensus implies vs the same quarter last year. Falling
        # back to last-reported YoY would mislabel the column.
        # Use same-source pairing: MS forecast vs MS prior-year actual.
        if next_est and prior_of_next:
            # ALL four metrics computed from MS pair when both sides exist.
            # Previously revenue + EPS deferred to Investing — but if MS
            # already has the forecast AND the prior-year actual, that's
            # the cleanest YoY: same source, same units, no scale fights.
            yoy_rev    = _yoy_pct(next_est.get("revenue"),    prior_of_next.get("revenue"))
            yoy_ebitda = _yoy_pct(next_est.get("ebitda"),     prior_of_next.get("ebitda"))
            yoy_nii    = _yoy_pct(next_est.get("nii"),        prior_of_next.get("nii"))
            yoy_ni     = _yoy_pct(next_est.get("net_income"), prior_of_next.get("net_income"))
            yoy_eps    = _yoy_pct(next_est.get("eps"),        prior_of_next.get("eps"))
            curr_rev_ms = next_est.get("revenue") or 0
            prev_rev_ms = prior_of_next.get("revenue") or 0
            curr_eb_ms  = next_est.get("ebitda")
            prev_eb_ms  = prior_of_next.get("ebitda")
            curr_margin = (curr_eb_ms / curr_rev_ms * 100) if curr_eb_ms and curr_rev_ms else None
            prev_margin = (prev_eb_ms / prev_rev_ms * 100) if prev_eb_ms and prev_rev_ms else None
            yoy_margin = _yoy_bps(curr_margin, prev_margin)
            used_ms = True
        # QoQ: next-Q forecast vs immediately-prior-quarter actual.
        if next_est and prior_quarter:
            qoq_rev    = _yoy_pct(next_est.get("revenue"),    prior_quarter.get("revenue"))
            qoq_eps    = _yoy_pct(next_est.get("eps"),        prior_quarter.get("eps"))
            qoq_ebitda = _yoy_pct(next_est.get("ebitda"),     prior_quarter.get("ebitda"))
            qoq_nii    = _yoy_pct(next_est.get("nii"),        prior_quarter.get("nii"))
            qoq_ni     = _yoy_pct(next_est.get("net_income"), prior_quarter.get("net_income"))
            curr_rev_q = next_est.get("revenue") or 0
            prev_rev_q = prior_quarter.get("revenue") or 0
            curr_eb_q  = next_est.get("ebitda")
            prev_eb_q  = prior_quarter.get("ebitda")
            curr_margin_q = (curr_eb_q / curr_rev_q * 100) if curr_eb_q and curr_rev_q else None
            prev_margin_q = (prev_eb_q / prev_rev_q * 100) if prev_eb_q and prev_rev_q else None
            qoq_margin = _yoy_bps(curr_margin_q, prev_margin_q)
        if not (next_est and prior_of_next) and latest and prior:
            # No forecast available for the next quarter (e.g. BKMB has no
            # MS Q2 forecast). Leave YoY blank rather than show a misleading
            # last-reported YoY in a forecast-labeled table.
            used_ms = True

    # Memo-cascade fallback for Q+1 consensus — per-field, not all-or-
    # nothing. The previous "fire only when both rev AND eps are None"
    # short-circuit meant that a stale `revenue_next_q` from canonical
    # store (e.g. BKMB's Q1'26 forecast still in MS summary, served as
    # the Q2'26 estimate by mistake) blocked MS-annual÷4 from filling
    # the OTHER columns (Net Income, EPS). Now each field gets its own
    # None-check, so the cascade can supplement rather than only
    # replace. `_compute_memo` in build_report_payload has already
    # walked MS quarterly → Investing → Yahoo → MS annual÷4; whatever
    # it found is on memo_data here.
    if memo_data:
        mc_rev = memo_data.get("next_quarter_consensus_revenue")
        mc_eps = memo_data.get("next_quarter_consensus_eps")
        mc_eb  = memo_data.get("next_quarter_consensus_ebitda")
        mc_ni  = memo_data.get("next_quarter_consensus_ni")
        if rev_q_consensus is None and isinstance(mc_rev, (int, float)):
            rev_q_consensus = mc_rev
        if eps_q_consensus is None and isinstance(mc_eps, (int, float)):
            eps_q_consensus = mc_eps
        if ebitda_q_consensus is None and isinstance(mc_eb, (int, float)):
            ebitda_q_consensus = mc_eb
        if ni_q_consensus is None and isinstance(mc_ni, (int, float)):
            ni_q_consensus = mc_ni

    # ── Yahoo backbone fallback for the CONSENSUS column ────────────────
    # When neither MS quarterly forecasts nor the memo cascade produced a
    # next-quarter consensus (the common case for Yahoo-covered EM names
    # with no MS slug — SNB, SMIC, Ooredoo …), fall back to Yahoo's analyst
    # estimate panel, already in canonical `valuation_forward`:
    #   eps_next_q     → EPS row consensus
    #   revenue_next_q → Revenue / (bank) Operating Income row consensus
    #   net income     → derived as eps_next_q × shares (shares = mcap/price)
    # This is what keeps the slide-2 table from shipping half-empty.
    if eps_q_consensus is None and isinstance(fwd.get("eps_next_q"), (int, float)):
        eps_q_consensus = fwd.get("eps_next_q")
    if rev_q_consensus is None and isinstance(fwd.get("revenue_next_q"), (int, float)):
        rev_q_consensus = fwd.get("revenue_next_q")
    if ni_q_consensus is None and isinstance(eps_q_consensus, (int, float)):
        try:
            _mc = cv.get("market_cap"); _cp = cv.get("current_price")
            _mcv = _mc.value if _mc else None
            _cpv = _cp.value if _cp else None
            if isinstance(_mcv, (int, float)) and isinstance(_cpv, (int, float)) and _cpv > 0:
                _shares = _mcv / _cpv
                ni_q_consensus = eps_q_consensus * _shares
        except Exception:
            pass

    # YoY / QoQ for tickers where `_ms_quarterly_split` returned empty
    # `next_est` (all MS dates already past → all rows bucketed as
    # actuals → no estimate row to anchor against). For BKMB-shape
    # names, the Q+1 forecast is now in rev_q_consensus / eps_q_consensus
    # / ni_q_consensus / ebitda_q_consensus (from the memo cascade
    # above), and the prior-year-same-Q and immediately-prior-Q actuals
    # are in MS quarterly periods data. Look them up directly here
    # rather than going through the past/future bucketing.
    if (ms_quarterly_forecasts and isinstance(ms_quarterly_forecasts, dict)
        and (yoy_rev is None or yoy_ni is None
             or qoq_rev is None or qoq_ni is None)):
        nq_label = (memo_data or {}).get("next_quarter_label") or ""
        import re as _re_q
        _m_nq = _re_q.search(
            r"(\d{4})\s*Q([1-4])|Q([1-4])\s*(\d{4})", nq_label
        )
        if _m_nq:
            _nq_yr = int(_m_nq.group(1) or _m_nq.group(4))
            _nq_q  = int(_m_nq.group(2) or _m_nq.group(3))
            _prev_q_q  = _nq_q - 1 if _nq_q > 1 else 4
            _prev_q_yr = _nq_yr if _nq_q > 1 else _nq_yr - 1

            _qtr = ms_quarterly_forecasts.get("quarterly", {}) or {}
            _periods = _qtr.get("periods", []) or []

            def _find_period_idx(yr, q):
                """Find the index in MS periods array for {yr}Q{q}."""
                want = f"{yr}Q{q}"
                for i, p in enumerate(_periods):
                    if str(p or "").replace(" ", "") == want:
                        return i
                return None

            def _raw_at(arr_name, idx):
                """Return the raw MS quarterly value (no unit scaling)."""
                if idx is None: return None
                arr = _qtr.get(arr_name, []) or []
                if idx >= len(arr): return None
                v = arr[idx]
                if not isinstance(v, (int, float)): return None
                return float(v)

            # Bring the consensus value to RAW MS scale so percent change
            # is computed against same-units denominators. Two paths fill
            # rev_q_consensus etc.:
            #   (a) `_scale(next_est.get(...))` inside `used_ms` — already
            #       multiplied by unit_scale (absolute currency units).
            #   (b) `fwd.get("revenue_next_q")` / memo cascade — raw, in
            #       the MS reported unit (millions).
            # Percent change is scale-invariant ONLY when both sides match;
            # mixing absolute vs millions yielded the -100% YoY/QoQ bug.
            def _to_raw_ms(v):
                if not isinstance(v, (int, float)): return None
                return v / unit_scale if used_ms else v

            _raw_rev_cons = _to_raw_ms(rev_q_consensus)
            _raw_ni_cons  = _to_raw_ms(ni_q_consensus)
            _raw_eb_cons  = _to_raw_ms(ebitda_q_consensus)
            # EPS is per-share — never unit-scaled on either side.

            _py_idx = _find_period_idx(_nq_yr - 1, _nq_q)
            _pq_idx = _find_period_idx(_prev_q_yr, _prev_q_q)

            # Per-field YoY (vs prior-year same Q actual)
            def _pct(num, den):
                if not (isinstance(num, (int, float))
                        and isinstance(den, (int, float)) and den != 0):
                    return None
                return (num - den) / den * 100.0

            if yoy_rev is None:
                yoy_rev = _pct(_raw_rev_cons, _raw_at("net_sales",  _py_idx))
            if yoy_ni is None:
                yoy_ni  = _pct(_raw_ni_cons,  _raw_at("net_income", _py_idx))
            if yoy_ebitda is None:
                yoy_ebitda = _pct(_raw_eb_cons, _raw_at("ebitda",  _py_idx))
            if yoy_eps is None:
                _eps_arr = _qtr.get("eps", []) or []
                _py_eps  = _eps_arr[_py_idx] if (_py_idx is not None and _py_idx < len(_eps_arr)) else None
                yoy_eps = _pct(eps_q_consensus, _py_eps)

            # Per-field QoQ (vs immediately-prior Q actual)
            if qoq_rev is None:
                qoq_rev = _pct(_raw_rev_cons, _raw_at("net_sales",  _pq_idx))
            if qoq_ni is None:
                qoq_ni  = _pct(_raw_ni_cons,  _raw_at("net_income", _pq_idx))
            if qoq_ebitda is None:
                qoq_ebitda = _pct(_raw_eb_cons, _raw_at("ebitda",  _pq_idx))
            if qoq_eps is None:
                _eps_arr = _qtr.get("eps", []) or []
                _pq_eps  = _eps_arr[_pq_idx] if (_pq_idx is not None and _pq_idx < len(_eps_arr)) else None
                qoq_eps = _pct(eps_q_consensus, _pq_eps)

    # Yahoo quarterly_actuals fallback (used when MS forecast block was empty
    # or didn't yield a YoY pair).
    if not used_ms and quarterly:
        # Quarterly list is expected to be FinancialPeriod objects (or dicts);
        # render side accepts both since serialization paths differ.
        def _g(rec, key):
            return getattr(rec, key, None) if not isinstance(rec, dict) else rec.get(key)
        recs_by_period = {(_g(r, "period_label") or "").strip(): r for r in quarterly}
        sorted_periods = sorted(recs_by_period.keys(), reverse=True)
        latest = recs_by_period.get(sorted_periods[0]) if sorted_periods else None
        prior_key = _prior_year_same_q(sorted_periods[0]) if sorted_periods else None
        prior = recs_by_period.get(prior_key) if prior_key else None
        if latest and prior:
            yoy_rev    = _yoy_pct(_g(latest, "revenue"),    _g(prior, "revenue"))
            yoy_ebitda = _yoy_pct(_g(latest, "ebitda"),     _g(prior, "ebitda"))
            yoy_nii    = _yoy_pct(_g(latest, "nii"),        _g(prior, "nii"))
            yoy_ni     = _yoy_pct(_g(latest, "net_income"), _g(prior, "net_income"))
            yoy_eps    = _yoy_pct(_g(latest, "eps"),        _g(prior, "eps"))
            # Margin YoY is computed as a bps delta, but we render it as a
            # percent-point change to keep the column format uniform.
            curr_rev = _g(latest, "revenue") or 0
            prev_rev = _g(prior, "revenue") or 0
            curr_eb  = _g(latest, "ebitda")
            prev_eb  = _g(prior, "ebitda")
            curr_margin = (curr_eb / curr_rev * 100) if curr_eb and curr_rev else None
            prev_margin = (prev_eb / prev_rev * 100) if prev_eb and prev_rev else None
            yoy_margin = _yoy_bps(curr_margin, prev_margin)

    # Investing.com earnings page as last-resort YoY (revenue + eps only).
    # Required for yfinance-blocked tickers where MS forecast block lacks
    # historical actuals or the announcement-date split fails.
    if (yoy_rev is None or yoy_eps is None) and ticker:
        inv = _investing_actuals_yoy(ticker)
        if yoy_rev is None and isinstance(inv.get("revenue"), (int, float)):
            yoy_rev = inv["revenue"]
        if yoy_eps is None and isinstance(inv.get("eps"), (int, float)):
            yoy_eps = inv["eps"]

    # Pick a single magnitude unit for the table — keeps values
    # comparable across rows. Based on the largest absolute value
    # across Revenue / EBITDA / Net Income consensus.
    abs_vals = [v for v in (rev_q_consensus, ebitda_q_consensus, ni_q_consensus, nii_q_consensus)
                  if isinstance(v, (int, float))]
    max_abs = max((abs(v) for v in abs_vals), default=0)
    if   max_abs >= 1e12: unit_div, unit_tag = 1e12, "T"
    elif max_abs >= 1e9:  unit_div, unit_tag = 1e9,  "B"
    elif max_abs >= 1e6:  unit_div, unit_tag = 1e6,  "M"
    else:                 unit_div, unit_tag = 1.0,  ""

    cur = (currency or "").upper()
    unit_suffix = (f"{cur}{unit_tag}" if cur and unit_tag else (cur or unit_tag))
    if unit_tag == "T":   subtitle_units = f"{cur} trillions unless stated"
    elif unit_tag == "B": subtitle_units = f"{cur} billions unless stated"
    elif unit_tag == "M": subtitle_units = f"{cur} millions unless stated"
    else:                  subtitle_units = f"{cur} units unless stated".strip()

    def _money_in_unit(v):
        if not isinstance(v, (int, float)): return None
        scaled = v / unit_div
        # Display precision: 1,000-ish numbers as integer; smaller with 1 dp.
        if abs(scaled) >= 100:
            return f"{scaled:,.0f}"
        if abs(scaled) >= 10:
            return f"{scaled:,.1f}"
        return f"{scaled:,.2f}"

    def _eps_fmt(v):
        if not isinstance(v, (int, float)): return None
        return f"{v:,.2f}"

    def _margin_fmt(v):
        if not isinstance(v, (int, float)): return None
        return f"{v:.1f}%"

    def _row(metric: str, jabal_str: str | None, consensus_str: str | None,
              yoy_val, qoq_val, is_margin: bool = False) -> dict:
        return {
            "metric":    metric,
            "jabal":     jabal_str if jabal_str else "—",
            "yoy":       yoy_val if yoy_val is not None else "—",
            "qoq":       qoq_val if qoq_val is not None else "—",
            "consensus": consensus_str if consensus_str else "—",
            "is_margin": is_margin,
        }

    eps_consensus_str = _eps_fmt(eps_q_consensus)

    if is_bank:
        # Bank schema: Operating Income (NII + non-int) / Net Income / EPS.
        # NII / Non-Int / PPOP / Provisions need broker-level data that
        # MS and Investing don't separately publish — Bloomberg upload is
        # the right path for those rows when wanted.
        rows = [
            _row(f"Operating Income ({unit_suffix})",
                  None, _money_in_unit(rev_q_consensus), yoy_rev, qoq_rev),
            _row(f"Net Income ({unit_suffix})",
                  None, _money_in_unit(ni_q_consensus),  yoy_ni,  qoq_ni),
            _row(f"EPS ({cur})",
                  None, eps_consensus_str,               yoy_eps, qoq_eps),
        ]
        return rows, unit_suffix

    # Non-bank: Revenue / EBITDA / Net Income / EPS / EBITDA Margin.
    # EBITDA Margin = EBITDA / Revenue × 100, computed from forecast pair.
    margin_consensus_str = None
    if (isinstance(ebitda_q_consensus, (int, float))
        and isinstance(rev_q_consensus, (int, float))
        and rev_q_consensus > 0):
        margin_consensus_str = _margin_fmt(ebitda_q_consensus / rev_q_consensus * 100)

    rows = [
        _row(f"Revenue ({unit_suffix})",     None, _money_in_unit(rev_q_consensus),    yoy_rev,    qoq_rev),
        _row(f"EBITDA ({unit_suffix})",      None, _money_in_unit(ebitda_q_consensus), yoy_ebitda, qoq_ebitda),
        _row(f"Net Income ({unit_suffix})",  None, _money_in_unit(ni_q_consensus),     yoy_ni,     qoq_ni),
        _row(f"EPS ({cur})",                  None, eps_consensus_str,                  yoy_eps,    qoq_eps),
        _row("EBITDA Margin",                 None, margin_consensus_str,               yoy_margin, qoq_margin,
              is_margin=True),
    ]
    return rows, unit_suffix


def _build_annual_rows(*, ms_annual_forecasts: dict | None,
                          ms_eps_dividend_forecasts: dict | None,
                          is_bank: bool, currency: str,
                          bloomberg_bundle: dict | None = None,
                          ) -> tuple[list[dict], list[str], str]:
    """Build the annual-FY fallback table rows.

    Returns (rows, fy_labels, unit_suffix). The first FY is treated as the
    YoY base (last actual). The remaining 3 are forecasts FY+1E/+2E/+3E.

    Returns ([], [], '') when MS doesn't carry annual data — caller falls
    back to the quarterly path.

    SOURCE PRIORITY (Bloomberg overrides everything rule):
      1. `bloomberg_bundle.annuals` when an uploaded Bloomberg FA xlsx
         is on disk for the ticker — the deck shows the same FY values
         the analyst sees on their BBG screen.
      2. `ms_annual_forecasts` + `ms_eps_dividend_forecasts` — the
         MarketScreener /finances/ + /valuation-dividend/ blocks.
    """
    # ── Bloomberg path ─────────────────────────────────────────────
    bbg_annuals = []
    if isinstance(bloomberg_bundle, dict):
        bbg_annuals = list(bloomberg_bundle.get("annuals") or [])
    if bbg_annuals:
        # Strip Current/LTM rows so we don't double-count alongside the
        # chronological actuals + estimates.
        bbg_annuals = [a for a in bbg_annuals if not a.get("is_ltm")]
        # Anchor YoY against the LAST ACTUAL year, not "the entry before
        # the first forecast we picked." With BBG showing FY25 Act + FY26-
        # FY29 Est, the naive last-4 slice picked FY26 as the YoY base
        # and made the FY27E YoY column look like 0% growth. Correct:
        # base = last `is_estimate=False`, fwds = first 3 estimates after it.
        actuals = [a for a in bbg_annuals if not a.get("is_estimate")]
        estimates = [a for a in bbg_annuals if a.get("is_estimate")]
        if not actuals or not estimates:
            return [], [], ""
        base = actuals[-1]
        fwds = estimates[:3]
        unit_suffix = f"{currency.upper()}M" if currency else "M"

        def _money(v):
            if not isinstance(v, (int, float)): return None
            if abs(v) >= 100: return f"{v:,.0f}"
            if abs(v) >= 10:  return f"{v:,.1f}"
            return f"{v:,.2f}"

        def _eps_fmt(v):
            # 3 decimals — Bloomberg EPS values are sub-1.0 for many
            # banks (BKMB FY26E = 0.035). At 2 decimals the entire
            # forecast strip rounds to "0.03 / 0.04 / 0.04" which
            # hides the trajectory the IC reader is looking for.
            if not isinstance(v, (int, float)): return None
            return f"{v:,.3f}"

        def _yoy(num, den):
            if not (isinstance(num, (int, float))
                    and isinstance(den, (int, float)) and den != 0):
                return None
            return (num - den) / den * 100.0

        import re as _re_bbg_fy
        from datetime import datetime as _dt_bbg
        _current_year = _dt_bbg.utcnow().year
        def _fy_disp(label):
            """Suffix convention requested by Mohamed (2026-05):
                 A = actual (past FY)
                 E = expected (current FY — book closes this year)
                 F = forecast (future FYs beyond current)
               Bloomberg labels are "FY 2026 Est" / "FY 2025"; we use the
               year + is_estimate flag to assign A/E/F."""
            s = str(label or "")
            m = _re_bbg_fy.search(r"(\d{4})", s)
            if not m: return s
            yr = int(m.group(1))
            if "est" not in s.lower():
                return f"FY{yr}A"
            return f"FY{yr}E" if yr == _current_year else f"FY{yr}F"

        fy_disp = [_fy_disp(f.get("period_label")) for f in fwds]
        while len(fy_disp) < 3:
            fy_disp.append("—")

        base_m = base.get("metrics") or {}
        def _m(d, k):
            v = d.get(k)
            return float(v) if isinstance(v, (int, float)) else None

        if is_bank:
            metric_specs = [
                (f"Revenue ({unit_suffix})", "revenue",    _money),
                (f"Net Income ({unit_suffix})", "net_income", _money),
                (f"EPS ({currency})",       "eps",        _eps_fmt),
            ]
        else:
            metric_specs = [
                (f"Revenue ({unit_suffix})", "revenue",    _money),
                (f"EBITDA ({unit_suffix})", "ebitda",     _money),
                (f"Net Income ({unit_suffix})", "net_income", _money),
                (f"EPS ({currency})",       "eps",        _eps_fmt),
            ]

        def _cagr(end_v, base_v, years):
            """Annualised compound growth rate as a percentage."""
            if not (isinstance(end_v, (int, float))
                    and isinstance(base_v, (int, float))
                    and base_v > 0 and end_v > 0 and years > 0):
                return None
            return ((end_v / base_v) ** (1.0 / years) - 1.0) * 100.0

        rows: list[dict] = []
        for label, key, fmt in metric_specs:
            base_v = _m(base_m, key)
            fwd_vals = [_m(f.get("metrics") or {}, key) for f in fwds]
            n_fwd = len(fwd_vals)
            while len(fwd_vals) < 3:
                fwd_vals.append(None)
            yoy_pct = _yoy(fwd_vals[0], base_v)
            # CAGR uses the LAST non-None forecast as the endpoint and the
            # base actual as the start. Years = position of that endpoint
            # in the forecast strip (1, 2, or 3 depending on coverage).
            cagr_pct = None
            last_idx = None
            for i in range(min(n_fwd, 3) - 1, -1, -1):
                if fwd_vals[i] is not None:
                    last_idx = i; break
            if last_idx is not None:
                cagr_pct = _cagr(fwd_vals[last_idx], base_v, last_idx + 1)
            rows.append({
                "metric": label,
                "fy1": fmt(fwd_vals[0]) if fwd_vals[0] is not None else None,
                "fy2": fmt(fwd_vals[1]) if fwd_vals[1] is not None else None,
                "fy3": fmt(fwd_vals[2]) if fwd_vals[2] is not None else None,
                "yoy": yoy_pct,
                "cagr": cagr_pct,
            })
        if any(r["fy1"] is not None for r in rows):
            return rows, fy_disp, unit_suffix

    # ── MS path ────────────────────────────────────────────────────
    if not isinstance(ms_annual_forecasts, dict):
        return [], [], ""
    annual = ms_annual_forecasts.get("annual") or {}
    if not isinstance(annual, dict):
        return [], [], ""
    periods_all = list(annual.get("periods") or [])
    if not periods_all:
        return [], [], ""
    # MS publishes periods chronologically. Heuristic: the first index
    # whose label doesn't END with 'A' (some publishers mark actuals) and
    # whose year >= current year is the first forecast. Conservative
    # alternative: use the last period as FY+3E, prior 3 are the FY series.
    # We pick the LAST 4 periods so a stale snapshot still surfaces useful
    # forward years. The first of those 4 is treated as the YoY base.
    if len(periods_all) < 2:
        return [], [], ""
    take = periods_all[-4:] if len(periods_all) >= 4 else periods_all
    base_label = take[0]
    fy_labels_raw = take[1:]   # forecasts (1-3 entries)
    # Indices into the source arrays
    n = len(periods_all)
    base_idx = n - len(take)
    fwd_indices = list(range(base_idx + 1, n))

    def _arr_at(arr_name: str, idx: int):
        arr = annual.get(arr_name) or []
        if 0 <= idx < len(arr):
            v = arr[idx]
            return float(v) if isinstance(v, (int, float)) else None
        return None

    # EPS sits in the dividend/EPS payload, possibly flat or wrapped.
    eps_periods: list = []
    eps_arr: list = []
    if isinstance(ms_eps_dividend_forecasts, dict):
        ed_wrapped = ms_eps_dividend_forecasts.get("annual")
        if isinstance(ed_wrapped, dict) and ed_wrapped.get("periods"):
            eps_periods = list(ed_wrapped.get("periods") or [])
            eps_arr = list(ed_wrapped.get("eps") or [])
        else:
            eps_periods = list(ms_eps_dividend_forecasts.get("periods") or [])
            eps_arr = list(ms_eps_dividend_forecasts.get("eps") or [])

    def _eps_at(period_label):
        for i, p in enumerate(eps_periods):
            if str(p) == str(period_label) and i < len(eps_arr):
                v = eps_arr[i]
                return float(v) if isinstance(v, (int, float)) else None
        return None

    # Unit & currency formatting — mirror the quarterly table conventions.
    unit_scale_label = (ms_annual_forecasts.get("unit_scale") or "").strip().lower()
    if unit_scale_label.startswith("bil"):
        unit_suffix = f"{currency.upper()}B" if currency else "B"
        unit_div = 1.0   # values already in billions
    elif unit_scale_label.startswith("thou"):
        unit_suffix = f"{currency.upper()}M" if currency else "M"
        unit_div = 1000.0   # convert thousands → millions for display
    else:
        unit_suffix = f"{currency.upper()}M" if currency else "M"
        unit_div = 1.0   # MS default = millions

    def _money(v):
        if not isinstance(v, (int, float)): return None
        scaled = v / unit_div
        if abs(scaled) >= 100:  return f"{scaled:,.0f}"
        if abs(scaled) >= 10:   return f"{scaled:,.1f}"
        return f"{scaled:,.2f}"

    def _eps_fmt(v):
        # Adaptive precision: sub-1.0 EPS (most GCC banks — BKMB FY26E ≈
        # 0.035) needs 3 decimals, otherwise the whole forecast strip
        # rounds to "0.03 / 0.04 / 0.04" and hides the trajectory the IC
        # reader is looking for. Larger EPS keeps 2 decimals.
        if not isinstance(v, (int, float)): return None
        return f"{v:,.3f}" if abs(v) < 1 else f"{v:,.2f}"

    def _yoy(num, den):
        if not (isinstance(num, (int, float))
                and isinstance(den, (int, float)) and den != 0):
            return None
        return (num - den) / den * 100.0

    # A / E / F suffix convention (Mohamed 2026-05):
    #   A = past actual
    #   E = current FY (expected)
    #   F = future FY (forecast)
    import re as _re_ms_fy
    from datetime import datetime as _dt_ms_fy
    _current_year_ms = _dt_ms_fy.utcnow().year
    def _fy_disp(p):
        s = str(p or "")
        m = _re_ms_fy.search(r"(\d{4})", s)
        if not m: return s
        yr = int(m.group(1))
        # MS forecast periods are all forward; we can't distinguish A
        # from E/F here cheaply. We DO know current_year, so:
        if yr < _current_year_ms:  return f"FY{yr}A"
        if yr == _current_year_ms: return f"FY{yr}E"
        return f"FY{yr}F"

    fy_disp = [_fy_disp(p) for p in fy_labels_raw]

    # Pull metric series across the 3 forecast indices + the base.
    # Bank P&L: net sales (total banking income) → EBIT (operating profit)
    # → pre-tax (EBT) → net income → EPS. EBITDA is meaningless for banks
    # (MS publishes it as 0.00) so it's never listed; the all-zero filter
    # below also drops any of these MS rows that come back empty for a
    # given name, so we never render a "0 / 0 / 0" line.
    metric_specs = []
    if is_bank:
        # For a bank, MS "net sales" IS total operating income (NII +
        # non-interest income) — not "revenue" in the industrial sense.
        # Label it accordingly so the deck reads correctly to a bank analyst.
        metric_specs = [
            (f"Total income ({unit_suffix})", "net_sales", _money),
            (f"EBIT ({unit_suffix})", "ebit", _money),
            (f"Pre-tax ({unit_suffix})", "ebt", _money),
            (f"Net Income ({unit_suffix})", "net_income", _money),
            (f"EPS ({currency})", None, _eps_fmt),   # EPS handled separately
        ]
    else:
        metric_specs = [
            (f"Revenue ({unit_suffix})", "net_sales", _money),
            (f"EBITDA ({unit_suffix})", "ebitda", _money),
            (f"EBIT ({unit_suffix})", "ebit", _money),
            (f"Pre-tax ({unit_suffix})", "ebt", _money),
            (f"Net Income ({unit_suffix})", "net_income", _money),
            (f"EPS ({currency})", None, _eps_fmt),
        ]

    def _row_has_signal(base_v, fwd_vals) -> bool:
        """True when the row carries at least one meaningful (non-zero,
        non-None) value across base + forecasts. Filters MS's all-0.00
        rows (EBITDA for banks; EBT/EBIT for names MS doesn't model) so the
        deck never shows a dead '0 / 0 / 0' line."""
        for v in [base_v, *fwd_vals]:
            if isinstance(v, (int, float)) and abs(v) > 1e-9:
                return True
        return False

    rows: list[dict] = []
    for label, arr_name, fmt in metric_specs:
        if arr_name == "net_sales" and "EPS" in label:
            continue   # never happens; defensive
        if arr_name is None:
            # EPS row
            base_v = _eps_at(base_label)
            fwd_vals = [_eps_at(p) for p in fy_labels_raw]
        else:
            base_v = _arr_at(arr_name, base_idx)
            fwd_vals = [_arr_at(arr_name, i) for i in fwd_indices]
        # Pad to 3 forecast columns
        while len(fwd_vals) < 3:
            fwd_vals.append(None)
        # Drop rows MS returns as all-zero / all-missing (e.g. bank EBITDA,
        # or EBIT/EBT for names MS doesn't break out) — but only the
        # forecast strip matters for "is this row worth showing".
        if not _row_has_signal(None, fwd_vals):
            continue
        yoy_pct = _yoy(fwd_vals[0], base_v) if fwd_vals else None
        # CAGR: last non-None forecast vs base, annualised.
        cagr_pct = None
        for i in range(min(3, len(fwd_vals)) - 1, -1, -1):
            if fwd_vals[i] is not None and isinstance(base_v, (int, float)) and base_v > 0:
                try:
                    if fwd_vals[i] > 0:
                        cagr_pct = ((fwd_vals[i] / base_v) ** (1.0 / (i + 1)) - 1.0) * 100.0
                except (TypeError, ValueError, ZeroDivisionError):
                    cagr_pct = None
                break
        rows.append({
            "metric": label,
            "fy1": fmt(fwd_vals[0]) if fwd_vals[0] is not None else None,
            "fy2": fmt(fwd_vals[1]) if fwd_vals[1] is not None else None,
            "fy3": fmt(fwd_vals[2]) if fwd_vals[2] is not None else None,
            "yoy": yoy_pct,
            "cagr": cagr_pct,
        })

    # If every value is None, treat as no data.
    if all(r["fy1"] is None and r["fy2"] is None and r["fy3"] is None for r in rows):
        return [], [], ""

    # Pad FY labels to length 3 for the column headers.
    while len(fy_disp) < 3:
        fy_disp.append("—")
    return rows, fy_disp, unit_suffix


def _build_annual_rows_from_yahoo(cv: dict, *, is_bank: bool, currency: str,
                                    ticker: str) -> tuple[list[dict], list[str], str]:
    """Annual table from Yahoo's FY estimates (canonical valuation_forward) +
    the grounded base-year actual. Used when MS/Bloomberg annual forecasts are
    absent but Yahoo covers the name, so slide-2 shows a real FY1E/FY2F strip
    (Total income / Revenue, Net income, EPS) instead of the sparse 3-row
    per-quarter shell. Net income is scaled off EPS growth so its YoY stays
    consistent with the EPS YoY.
    """
    import re as _re
    vf = cv.get("valuation_forward")
    vf = (getattr(vf, "value", None) if vf is not None else None) or {}
    if not isinstance(vf, dict):
        vf = {}

    # Pull FY estimates straight from Yahoo (0y = current FY, +1y = next FY).
    # The canonical valuation_forward is often Investing-sourced and carries
    # only FY1, so go direct to get BOTH years for the strip; fall back to
    # canonical (FY1 only) when Yahoo is unavailable (e.g. Yahoo-blind names).
    eps1 = eps2 = rev1 = rev2 = None
    try:
        from src.providers._yf import yf
        yt = yf.Ticker(ticker)

        def _avg(df, period):
            try:
                if df is None or period not in df.index:
                    return None
                v = float(df.loc[period].get("avg"))
                return v if v == v else None
            except Exception:
                return None
        eps1, eps2 = _avg(yt.earnings_estimate, "0y"), _avg(yt.earnings_estimate, "+1y")
        rev1, rev2 = _avg(yt.revenue_estimate, "0y"), _avg(yt.revenue_estimate, "+1y")
    except Exception:
        pass
    if eps1 is None and rev1 is None:
        eps1, eps2 = vf.get("eps_fy1"), vf.get("eps_fy2")
        rev1, rev2 = vf.get("revenue_fy1"), vf.get("revenue_fy2")
    if eps1 is None and rev1 is None:
        return [], [], ""

    from src.services.disclosed_loader import load_disclosed
    fyh = (load_disclosed(ticker) or {}).get("fy_highlights") or {}
    m = _re.search(r"(\d{4})", str(fyh.get("period", "")))
    base_year = int(m.group(1)) if m else None
    fy1_year = vf.get("fy1_year") or (base_year + 1 if base_year else None)
    fy2_year = vf.get("fy2_year") or (base_year + 2 if base_year else None)
    fy1_label = f"FY{fy1_year}E" if fy1_year else "FY+1E"
    fy2_label = f"FY{fy2_year}F" if fy2_year else "FY+2F"

    def _mn(k):
        v = fyh.get(k)
        return v * 1e6 if isinstance(v, (int, float)) else None
    base_rev = _mn("revenue_mn") or _mn("total_income_mn")
    if base_rev is None and is_bank:               # banks: total income = NII + non-interest
        nii, non = _mn("nii_mn"), _mn("non_interest_income_mn")
        base_rev = (nii or 0) + (non or 0) or None
    base_ni = _mn("net_profit_mn")
    base_eps = fyh.get("eps") if isinstance(fyh.get("eps"), (int, float)) else None

    # China A-shares etc. expose REVENUE estimates but no EPS — imply EPS from
    # the grounded base scaled by revenue growth (constant-margin) so the table
    # isn't a single Revenue row.
    if eps1 is None and isinstance(rev1, (int, float)) and base_rev and base_eps:
        eps1 = base_eps * (rev1 / base_rev)
        if isinstance(rev2, (int, float)):
            eps2 = base_eps * (rev2 / base_rev)

    # Net income tracks EPS growth off the grounded base (consistent YoY), or
    # revenue growth when EPS itself was implied above.
    ni1 = ni2 = None
    if base_ni and base_eps and isinstance(eps1, (int, float)):
        ni1 = base_ni * (eps1 / base_eps)
        if isinstance(eps2, (int, float)):
            ni2 = base_ni * (eps2 / base_eps)
    elif base_ni and base_rev and isinstance(rev1, (int, float)):
        ni1 = base_ni * (rev1 / base_rev)
        if isinstance(rev2, (int, float)):
            ni2 = base_ni * (rev2 / base_rev)

    nums = [abs(x) for x in (rev1, rev2, ni1, ni2) if isinstance(x, (int, float))]
    biggest = max(nums) if nums else 0
    if biggest >= 1e9:
        div, suffix = 1e9, (f"{currency.upper()}B" if currency else "B")
    else:
        div, suffix = 1e6, (f"{currency.upper()}M" if currency else "M")

    def _money(v):
        if not isinstance(v, (int, float)):
            return None
        s = v / div
        return f"{s:,.0f}" if abs(s) >= 100 else (f"{s:,.1f}" if abs(s) >= 10 else f"{s:,.2f}")

    def _eps(v):
        return f"{v:,.3f}" if isinstance(v, (int, float)) else None

    def _yoy(est, base):
        if not (isinstance(est, (int, float)) and isinstance(base, (int, float)) and base):
            return None
        return (est / base - 1.0) * 100.0

    def _cagr(fy2v, base):
        if not (isinstance(fy2v, (int, float)) and isinstance(base, (int, float))
                and base > 0 and fy2v > 0):
            return None
        return ((fy2v / base) ** 0.5 - 1.0) * 100.0

    rows: list[dict] = []
    if isinstance(rev1, (int, float)):
        rows.append({"metric": f"{'Total income' if is_bank else 'Revenue'} ({suffix})",
                     "fy1": _money(rev1), "fy2": _money(rev2), "fy3": "—",
                     "yoy": _yoy(rev1, base_rev), "cagr": _cagr(rev2, base_rev)})
    if isinstance(ni1, (int, float)):
        rows.append({"metric": f"Net Income ({suffix})",
                     "fy1": _money(ni1), "fy2": _money(ni2), "fy3": "—",
                     "yoy": _yoy(ni1, base_ni), "cagr": _cagr(ni2, base_ni)})
    if isinstance(eps1, (int, float)):
        rows.append({"metric": f"EPS ({currency.upper()})" if currency else "EPS",
                     "fy1": _eps(eps1), "fy2": _eps(eps2), "fy3": "—",
                     "yoy": _yoy(eps1, base_eps), "cagr": _cagr(eps2, base_eps)})
    if not rows:
        return [], [], ""
    return rows, [fy1_label, fy2_label], suffix


def build_thesis_data(ticker: str, *, analyst_name: str = "Jabal Research",
                        gen_date: str = "",
                        catalysts: Optional[list[str]] = None,
                        risks: Optional[list[str]] = None,
                        watch_list: Optional[list[str]] = None,
                        quarterly: Optional[list] = None,
                        is_bank: bool = False,
                        ms_quarterly_forecasts: Optional[dict] = None,
                        ms_annual_forecasts: Optional[dict] = None,
                        ms_eps_dividend_forecasts: Optional[dict] = None,
                        bloomberg_bundle: Optional[dict] = None,
                        period_heading: Optional[str] = None,
                        memo_data: Optional[dict] = None,
                        ) -> ThesisData:
    cv = get_all_fields(ticker)
    commodities_obs = get_observations_by_provider(ticker, "commodities")
    macro_obs       = get_observations_by_provider(ticker, "macro")
    investing_obs   = get_observations_by_provider(ticker, "investing")

    # LLM prose priority:
    #   1. pptx_sections — the pipeline's draft_pptx_sections output (single
    #      Gemini call producing thesis + catalysts + risks + watch together,
    #      so the four blocks stay mutually consistent and non-duplicative).
    #      This is the richer ~110-word thesis; reuse it instead of firing a
    #      SECOND Gemini call here.
    #   2. generate_summary — standalone fallback (e.g. build_thesis_data
    #      called directly without a pipeline, or for back-compat).
    #   3. deterministic template — when the LLM is off/unavailable.
    sections = (memo_data or {}).get("pptx_sections") or {}
    llm = None
    if isinstance(sections, dict) and (sections.get("investment_thesis") or "").strip():
        llm = {
            "thesis_paragraph": sections.get("investment_thesis"),
            "catalysts": sections.get("catalysts") or [],
            "risks": sections.get("risks") or [],
            "watch_list": sections.get("what_to_watch") or [],
        }
    else:
        try:
            from src.services.llm_summary import generate_summary
            llm = generate_summary(ticker)
        except Exception:
            llm = None
    summary = (llm or {}).get("thesis_paragraph") or _template_exec_summary(
        cv, commodities_obs, macro_obs, ticker=ticker)
    # Look up listing currency from company_master so the table can label
    # values "Revenue (SARM)" / "(AEDM)" etc. Falls back to canonical
    # profile currency when DB lookup misses.
    deck_currency = ""
    try:
        from src.storage.db import load_company as _lc
        cm = _lc(ticker) or {}
        deck_currency = (cm.get("currency") or "").strip()
    except Exception:
        pass
    if not deck_currency:
        prof = cv.get("company_profile")
        if prof and isinstance(prof.value, dict):
            deck_currency = (prof.value.get("currency") or "").strip()
    rows, unit_suffix = _build_estimates_rows(cv, quarterly=quarterly, is_bank=is_bank,
                                                ms_quarterly_forecasts=ms_quarterly_forecasts,
                                                ticker=ticker, currency=deck_currency,
                                                memo_data=memo_data)

    # Compose the table footnote: lead with the next-Q anchor (date,
    # consensus source, analyst count) since that's the strongest data
    # point we have for the full panel.
    val_fwd = cv.get("valuation_forward")
    rs = cv.get("rating_split")
    fwd_dict = val_fwd.value if val_fwd and isinstance(val_fwd.value, dict) else {}
    rs_dict  = rs.value if rs and isinstance(rs.value, dict) else {}
    n_an = int(rs_dict.get("total", 0) or 0)
    nq_period = fwd_dict.get("next_q_period") or ""
    nq_date   = fwd_dict.get("next_q_report_date") or ""
    fwd_source = (val_fwd.canonical_source if val_fwd else "—").title()
    # If the cascade fallback (`_compute_memo`) supplied Q+1 numbers because
    # MS / canonical_store was empty, credit that source instead of "—".
    memo_fwd_src = (memo_data or {}).get("next_quarter_consensus_source")
    if (not val_fwd or fwd_source in ("—", "")) and memo_fwd_src:
        fwd_source = memo_fwd_src
    footnote_bits = []
    if nq_period and nq_date:
        footnote_bits.append(f"Next print: {nq_date} (period {nq_period})")
    footnote_bits.append(f"Consensus: {fwd_source}")
    if n_an:
        footnote_bits.append(f"{n_an} analysts covering")
    estimates_footnote = "  ·  ".join(footnote_bits)

    # Surface Investing's surprise history as a track-record catalyst line.
    # Only when we actually have a non-trivial surprise — surfacing
    # "EPS missed by 0.0%; 0 of last 4 above estimates" on tickers with
    # missing data is worse than no line at all.
    surprise = (investing_obs.get("income_statement_quarterly") or {}).get(
        "surprise_history", [])
    track_record_catalyst = None
    if surprise:
        usable = [r for r in surprise[:4]
                  if isinstance(r.get("eps_surprise_pct"), (int, float))
                  and abs(r["eps_surprise_pct"]) >= 0.05]   # filter out 0/null rows
        beats = sum(1 for r in usable if r["eps_surprise_pct"] > 0)
        n = len(usable)
        last = next((r for r in surprise[:4]
                       if isinstance(r.get("eps_surprise_pct"), (int, float))
                       and abs(r["eps_surprise_pct"]) >= 0.05), None)
        if last and n > 0:
            last_dir = "beat" if last["eps_surprise_pct"] > 0 else "missed"
            last_pct = abs(last["eps_surprise_pct"])
            track_record_catalyst = (
                f"EPS {last_dir} consensus by {last_pct:.1f}% last quarter; "
                f"{beats} of last {n} quarters above estimates"
            )
    # Default catalyst / risk / watch lists. Sector-aware — banks
    # shouldn't fall back to "feedstock volatility" and chemicals
    # shouldn't fall back to "NIM trajectory". These templates only
    # ship when Gemini is unavailable; treat them as readable defaults
    # the analyst can rewrite, not analytical claims.
    _profile = cv.get("company_profile")
    _sector_l = ""
    _industry_l = ""
    if _profile and isinstance(_profile.value, dict):
        _sector_l = (_profile.value.get("sector") or "").lower()
        _industry_l = (_profile.value.get("industry") or "").lower()
    _is_bank_template = ("bank" in _sector_l or "bank" in _industry_l or "financial" in _sector_l)
    _is_energy_template = ("oil" in _industry_l or "gas" in _industry_l or "energy" in _sector_l)

    if _is_bank_template:
        default_catalysts_base = [
            "Net interest income trajectory vs prior quarter and management's NIM outlook",
            "Loan growth and deposit cost commentary into H2",
            "Capital deployment update — dividend cadence or buyback announcement",
        ]
        default_risks_template = [
            "Asset-quality deterioration / rising cost of risk pressuring earnings",
            "NIM compression as funding costs catch up with the rate cycle",
            "Loan-growth slowdown if domestic demand softens",
        ]
        default_watch_template = [
            "What is the trajectory of NIM and where does management see it stabilising?",
            "How is asset quality trending across the loan book?",
            "Any update on capital return policy through year-end?",
        ]
    elif _is_energy_template:
        default_catalysts_base = [
            "Production volume update and any guidance change",
            "Realized price commentary versus benchmark spot",
            "Capex / project-ramp timeline progress",
        ]
        default_risks_template = [
            "Commodity-price softness pressuring realized prices",
            "Project ramp slippage or cost overrun",
            "Capex intensity pressuring near-term free cash flow",
        ]
        default_watch_template = [
            "How is production tracking vs the full-year guidance?",
            "What is management's tone on realized-price discipline?",
            "Any commentary on the next capex / project milestone?",
        ]
    else:
        default_catalysts_base = [
            "Forward demand commentary and guidance update on the call",
            "Margin / cost trajectory versus prior quarter",
            "Capital return cadence (dividend / buyback) into year-end",
        ]
        default_risks_template = [
            "Cautious management tone could validate target-price gap",
            "Input / cost volatility pressures margin trajectory",
            "Macro / sector softness weighs on top-line growth",
        ]
        default_watch_template = [
            "Forward demand commentary — Q3 order book and pricing trajectory",
            "Cost outlook and supply-chain commentary",
            "Updated capex schedule and any project-pipeline updates",
        ]

    # Track-record line is included only when it carries real signal;
    # otherwise we lead with the sector-aware default.
    default_catalysts = [track_record_catalyst] + default_catalysts_base if track_record_catalyst else default_catalysts_base
    default_catalysts = default_catalysts[:3]

    # LLM output, when present, replaces every default bullet list too —
    # otherwise the deck mixes a fresh LLM thesis with stale boilerplate.
    def _dict_to_str(d: dict) -> str:
        return (d.get("catalyst") or d.get("risk") or d.get("question")
                or d.get("watch") or d.get("text") or d.get("value") or d.get("bullet")
                or next((v for v in d.values() if isinstance(v, str)), ""))

    def _norm_bullets(items) -> list[str]:
        """Coerce bullet items to plain strings. Upstreams hand us a list[str]
        (Gemini thesis), a list[dict] like {'catalyst': '...'} / {'bullet': '...'},
        OR — from the cached-draft path — a STRING that already contains a dict
        repr ("{'bullet': '...'}"). All three must collapse to the inner text;
        otherwise the deck renders the literal dict repr on the slide."""
        import ast as _ast
        out = []
        for it in (items or []):
            s = ""
            if isinstance(it, dict):
                s = _dict_to_str(it)
            elif isinstance(it, str):
                t = it.strip()
                # Stringified dict? Parse it and pull the inner value.
                if (t.startswith("{") and t.endswith("}")
                        and ("'" in t or '"' in t) and ":" in t):
                    try:
                        parsed = _ast.literal_eval(t)
                        s = _dict_to_str(parsed) if isinstance(parsed, dict) else t
                    except (ValueError, SyntaxError):
                        s = t
                else:
                    s = it
            else:
                s = str(it)
            s = (s or "").strip()
            if s:
                out.append(s)
        return out
    llm_catalysts = _norm_bullets((llm or {}).get("catalysts"))
    llm_risks     = _norm_bullets((llm or {}).get("risks"))
    llm_watch     = _norm_bullets((llm or {}).get("watch_list"))
    # Caller-supplied lists may also arrive as dicts — normalize them too so
    # neither return site (annual / quarterly) can emit a dict repr.
    catalysts = _norm_bullets(catalysts) if catalysts else catalysts
    risks      = _norm_bullets(risks) if risks else risks
    watch_list = _norm_bullets(watch_list) if watch_list else watch_list

    # Derive the Jabal-estimate column header (e.g. "Q2 2026E") from the
    # period_heading. Handles both "Q2 2026 Earnings Expectations" and the
    # quarter-first "1Q26"/"1Q 2026" forms used by some names. When no quarter
    # can be parsed, fall back to a clean "ESTIMATE" — NEVER the section title
    # (which produced the "EARNINGS EXPECTATIONS" column-header bug on 2010.SR).
    estimates_period_label = "ESTIMATE"
    if period_heading:
        import re as _re_ph
        ph = period_heading.strip()
        m = _re_ph.search(r"\bQ([1-4])\s*'?(\d{2,4})\b", ph)        # Q2 2026 / Q2 '26
        m2 = _re_ph.search(r"\b([1-4])Q\s*'?(\d{2,4})\b", ph)        # 1Q26 / 1Q 2026
        hit = m or m2
        if hit:
            q = hit.group(1)
            yr = hit.group(2)
            yr = ("20" + yr) if len(yr) == 2 else yr
            estimates_period_label = f"Q{q} {yr}E"
    # Subtitle line under the section heading.
    subtitle_unit_phrase = (
        f"{(deck_currency or '').upper()} {('trillions' if unit_suffix.endswith('T') else 'billions' if unit_suffix.endswith('B') else 'millions' if unit_suffix.endswith('M') else 'units')} unless stated".strip()
    )
    estimates_subtitle = f"Jabal estimates vs. consensus  ·  {subtitle_unit_phrase}"

    # Footnote: source + analyst count + "Bps = basis points" disclosure when
    # margin row is present.
    consensus_source = fwd_source if fwd_source and fwd_source != "—" else "MarketScreener"
    footnote_bits = [f"Estimates: Jabal Research", f"Consensus: {consensus_source}"]
    if n_an:
        footnote_bits[-1] = f"Consensus: {consensus_source} ({n_an} analysts)"
    if not is_bank:
        footnote_bits.append("Bps = basis points")
    estimates_footnote = "  ·  ".join(footnote_bits)

    # ANALYST-VALIDITY GATE — three independent signals can swap the per-
    # quarter table for an annual FY view. Each addresses a real failure
    # mode observed in the wild:
    #
    #   (1) memo cascade source contains "annual ÷" — the upstream code
    #       explicitly flagged it used the synthetic ÷4 proxy.
    #   (2) MS quarterly section has no forward-dated rows — i.e. MS
    #       publishes annual forecasts but no per-quarter breakdown.
    #       This is the BKMB shape: the calendar path can still inject a
    #       stale prior-quarter actual as the "next Q" estimate, which
    #       overwrites the source label and hides (1).
    #   (3) Q+1 consensus equals the immediately-prior-quarter actual to
    #       within rounding — the carry-forward fingerprint. When the
    #       calendar tier-2 fallback misclassifies a Q1 actual as a Q2
    #       forecast, this is the only signal that catches it.
    #
    # ANY one of these is sufficient to swap to the annual table.
    _memo_src = (memo_data or {}).get("next_quarter_consensus_source") or ""
    _signal_explicit_annual = "annual" in _memo_src.lower() and "÷" in _memo_src

    _signal_no_quarterly_fwd = False
    if isinstance(ms_quarterly_forecasts, dict):
        try:
            _nx, _lt, _pr, _pon, _pq = _ms_quarterly_split(ms_quarterly_forecasts)
            if not _nx:
                # No forward-dated quarterly row at all — MS doesn't
                # publish a per-quarter forecast for this name.
                _signal_no_quarterly_fwd = True
        except Exception:
            _signal_no_quarterly_fwd = False
    else:
        # No quarterly payload at all → annual is the only honest view.
        _signal_no_quarterly_fwd = True

    # (3) — compare memo consensus to the immediately-prior-quarter
    # actual from MS quarterly arrays. If revenue + NI both match the
    # last actual to within 1% AND there's no actual newer than that,
    # this is the carry-forward bug.
    _signal_carry_forward = False
    if isinstance(ms_quarterly_forecasts, dict) and memo_data:
        try:
            _qq = (ms_quarterly_forecasts.get("quarterly") or {})
            _rev_arr = _qq.get("net_sales") or []
            _ni_arr  = _qq.get("net_income") or []
            # Find last non-None actual index (i.e. the most-recent value).
            _last_rev = next((v for v in reversed(_rev_arr)
                                if isinstance(v, (int, float))), None)
            _last_ni  = next((v for v in reversed(_ni_arr)
                                if isinstance(v, (int, float))), None)
            _mc_rev = memo_data.get("next_quarter_consensus_revenue")
            _mc_ni  = memo_data.get("next_quarter_consensus_ni") \
                       or memo_data.get("next_quarter_consensus_net_income")
            def _close(a, b, tol=0.01):
                if not (isinstance(a, (int, float))
                        and isinstance(b, (int, float)) and b != 0):
                    return False
                return abs(a - b) / abs(b) < tol
            # Carry-forward when BOTH metrics match — single-metric match
            # could legitimately be a flat consensus.
            if _close(_mc_rev, _last_rev) and _close(_mc_ni, _last_ni):
                _signal_carry_forward = True
        except Exception:
            pass

    # Bloomberg bundle on disk also triggers annual mode — if BBG is
    # available, the deck uses the analyst's screen values regardless of
    # whether MS has a per-quarter forecast.
    _signal_bloomberg = bool(bloomberg_bundle
                              and isinstance(bloomberg_bundle, dict)
                              and (bloomberg_bundle.get("annuals") or []))
    _use_annual = (_signal_explicit_annual
                    or _signal_no_quarterly_fwd
                    or _signal_carry_forward
                    or _signal_bloomberg)
    annual_rows: list[dict] = []
    annual_fy_labels: list[str] = []
    annual_unit_suffix = ""
    _annual_from_yahoo = False
    if _use_annual and (ms_annual_forecasts or bloomberg_bundle):
        annual_rows, annual_fy_labels, annual_unit_suffix = _build_annual_rows(
            ms_annual_forecasts=ms_annual_forecasts,
            ms_eps_dividend_forecasts=ms_eps_dividend_forecasts,
            bloomberg_bundle=bloomberg_bundle,
            is_bank=is_bank, currency=deck_currency,
        )
    if not annual_rows and _use_annual:
        # No MS/Bloomberg annual data, but the per-quarter table would be sparse
        # (no quarterly forecast). Build the FY strip from Yahoo's annual
        # estimates so the slide isn't a 3-row shell — Total income / Net income
        # / EPS across FY+1E/FY+2F with YoY vs the grounded base year.
        annual_rows, annual_fy_labels, annual_unit_suffix = _build_annual_rows_from_yahoo(
            cv, is_bank=is_bank, currency=deck_currency, ticker=ticker)
        _annual_from_yahoo = bool(annual_rows)
    if annual_rows:
        # Re-label the section heading and subtitle for annual mode.
        period_heading_final = "Annual Earnings Expectations"
        subtitle_unit_phrase_a = (
            f"{(deck_currency or '').upper()} {('trillions' if annual_unit_suffix.endswith('T') else 'billions' if annual_unit_suffix.endswith('B') else 'millions' if annual_unit_suffix.endswith('M') else 'units')} unless stated".strip()
        )
        _consensus_source_a = ("Bloomberg consensus" if _signal_bloomberg
                                  else "Yahoo Finance annual" if _annual_from_yahoo
                                  else "MarketScreener annual")
        estimates_subtitle_a = (
            f"Annual analyst consensus  ·  {subtitle_unit_phrase_a}  ·  "
            f"per-quarter breakdown unavailable for this name"
        )
        footnote_a = "  ·  ".join([
            "Estimates: Jabal Research",
            f"Consensus: {_consensus_source_a}",
            # The Yahoo annual path derives net income (and, for revenue-only
            # names, EPS) from consensus scaled by the grounded base, so be
            # explicit rather than implying every line is a separate forecast.
            ("Net income / EPS derived from consensus + last-FY margin"
             if _annual_from_yahoo else "YoY computed vs prior-FY actual"),
        ])
        from datetime import datetime
        return ThesisData(
            exec_summary_body=summary,
            estimates_rows=rows,
            estimates_footnote=footnote_a,
            estimates_subtitle=estimates_subtitle_a,
            estimates_period_label=estimates_period_label,
            catalysts=_norm_bullets(catalysts or llm_catalysts or default_catalysts),
            risks=_norm_bullets(risks or llm_risks or default_risks_template),
            watch_list=_norm_bullets(watch_list or llm_watch or default_watch_template),
            sources_line=_sources_line_from_cv(cv),
            analyst_name=analyst_name,
            gen_date=gen_date or datetime.utcnow().strftime("%d %b %Y"),
            period_heading=period_heading_final,
            is_annual_table=True,
            annual_rows=annual_rows,
            annual_fy_labels=annual_fy_labels,
        )

    from datetime import datetime
    return ThesisData(
        exec_summary_body=summary,
        estimates_rows=rows,
        estimates_footnote=estimates_footnote,
        estimates_subtitle=estimates_subtitle,
        estimates_period_label=estimates_period_label,
        catalysts=catalysts or llm_catalysts or default_catalysts,
        risks=risks or llm_risks or default_risks_template,
        watch_list=watch_list or llm_watch or default_watch_template,
        sources_line=_sources_line_from_cv(cv),
        analyst_name=analyst_name,
        gen_date=gen_date or datetime.utcnow().strftime("%d %b %Y"),
        period_heading=(period_heading or "Earnings Expectations"),
    )


def _sources_line_from_cv(cv: dict) -> str:
    """Re-export of render_jabal_snapshot._sources_line so all three
    slide builders attribute the same way (every contributing provider,
    not just per-field winners)."""
    from src.services.render_jabal_snapshot import _sources_line
    return _sources_line(cv)
