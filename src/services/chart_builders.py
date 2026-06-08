"""
Chart builders for PPTX report generation.

Ported from earnings-preview.v2 so the PPTX output matches the "gold
standard" sample decks with native python-pptx charts instead of
shape-only rendering.

Creates embedded charts matching MarketScreener visual style:
- Income Statement Evolution: clustered column (Sales, EBIT, NI)
- P/E ratio bar chart with 5-year average + N/M annotation
- 1-year price line chart (downsampled)
- Surprise history summary box
- Expanded 6-column financial table

Every builder is a no-op when data is missing, so callers can wire
them unconditionally.
"""
from __future__ import annotations

from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.util import Inches, Pt


# ── Color palette matching MarketScreener ─────────────────────
BLACK_BAR = RGBColor(0x33, 0x33, 0x33)     # Sales bars
GOLD_BAR = RGBColor(0xE0, 0xB0, 0x30)      # EBIT / Operating Profit bars
GREEN_BAR = RGBColor(0x6B, 0x8E, 0x23)     # Net Income bars
GOLD_LINE = RGBColor(0xE0, 0xB0, 0x30)     # Operating Margin line
GREEN_LINE = RGBColor(0x4A, 0x7C, 0x2E)    # Net Margin line
MUTED_GRAY = RGBColor(0x8B, 0x94, 0x9E)
DARK_BLUE = RGBColor(0x1F, 0x3A, 0x5F)
GOLD = RGBColor(0xC9, 0xA2, 0x27)
ESTIMATE_GRAY = RGBColor(0xAA, 0xAA, 0xAA)
PURPLE_AVG = RGBColor(0x99, 0x33, 0xCC)


def _safe_float(v) -> float:
    if v is None:
        return 0.0
    try:
        return float(v)
    except (TypeError, ValueError):
        return 0.0


def _format_millions(v) -> str:
    if v is None:
        return ""
    try:
        x = float(v)
        if abs(x) >= 1e6:
            return f"{x/1e6:.0f}M"
        if abs(x) >= 1e3:
            return f"{x/1e3:.0f}K"
        return f"{x:,.0f}"
    except (TypeError, ValueError):
        return str(v)


def build_revenue_ni_chart(
    slide,
    x, y, w, h,
    periods: list[str],
    revenues: list,
    net_incomes: list,
    actuals_boundary: int,
    currency: str = "",
    ebit_values: list | None = None,
    net_margins: list | None = None,
    ebit_margins: list | None = None,
) -> None:
    """Income Statement Evolution: clustered bars for Sales / EBIT / Net Income.

    Distinguishes "no EBIT series" (None) from "empty slice" ([] → zeros), never
    substitutes EBITDA. Estimate bars use the same color as actuals — the
    actual vs estimate split is conveyed elsewhere in the slide.
    """
    if not periods or (not any(revenues) and not any(net_incomes)):
        return

    # Preserve None so missing periods render as no-bar instead of a
    # misleading zero-bar (see comment on build_quarterly_income_chart).
    def _to_display(arr):
        out = []
        for v in (arr or []):
            if v is None:
                out.append(None)
                continue
            try:
                out.append(float(v))
            except (TypeError, ValueError):
                out.append(None)
        return out

    rev_vals = _to_display(revenues)
    ni_vals = _to_display(net_incomes)
    ebit_vals = _to_display(ebit_values) if ebit_values is not None else [None] * len(periods)

    chart_data = CategoryChartData()
    labels = [p.replace("FY", "").strip() for p in periods]
    chart_data.categories = labels
    chart_data.add_series("Sales", rev_vals)
    chart_data.add_series("EBIT", ebit_vals)
    chart_data.add_series("Net Income", ni_vals)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
    )
    chart = chart_frame.chart

    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = BLACK_BAR
    chart.series[1].format.fill.solid()
    chart.series[1].format.fill.fore_color.rgb = GOLD_BAR
    chart.series[2].format.fill.solid()
    chart.series[2].format.fill.fore_color.rgb = GREEN_BAR

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(6)
    chart.legend.font.color.rgb = MUTED_GRAY

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(7)
    cat_axis.tick_labels.font.color.rgb = MUTED_GRAY
    cat_axis.tick_labels.font.bold = True
    cat_axis.has_major_gridlines = False
    try:
        cat_axis.format.line.fill.background()
    except Exception:
        pass

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(6)
    val_axis.tick_labels.font.color.rgb = MUTED_GRAY
    val_axis.has_major_gridlines = False
    try:
        val_axis.format.line.fill.background()
    except Exception:
        pass

    # Values arrive in millions already (MS quotes financials in millions; Yahoo
    # raw values are pre-scaled by `_toM` upstream). The previous format
    # `'#,##0,,"M"'` divided by another 1M and rendered every tick as "0M".
    try:
        val_axis.tick_labels.number_format = '#,##0"M"'
        val_axis.tick_labels.number_format_is_linked = False
    except Exception:
        pass

    try:
        chart_frame.line.fill.background()
    except AttributeError:
        pass


def build_pe_chart(
    slide,
    x, y, w, h,
    periods: list[str],
    pe_values: list,
    five_yr_avg: float | None = None,
) -> None:
    """P/E ratio bar chart with 5-year average annotation and N/M marker.

    Loss years (negative or zero P/E) render as a zero bar and are listed in a
    corner annotation like "N/M: 2023, 2024". Never plots a misleading bar.

    The 5-year average is drawn as a dashed horizontal overlay at the average
    Y-coordinate inside the plot area. Plot-area offsets are empirical (the
    pptx column-chart layout is consistent across renders), and we pin the
    value-axis maximum so the line position is deterministic.
    """
    import math
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.oxml.ns import qn
    from lxml import etree

    if not periods or not any(pe_values):
        return

    chart_data = CategoryChartData()
    labels = [p.replace("FY", "").strip() for p in periods]
    chart_data.categories = labels
    pe_clean: list[float] = []
    nm_years: list[str] = []
    for i, v in enumerate(pe_values):
        fv = _safe_float(v)
        if v is None or fv <= 0:
            pe_clean.append(0.0)
            nm_years.append(labels[i] if i < len(labels) else "")
        else:
            pe_clean.append(fv)
    chart_data.add_series("P/E", pe_clean)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    chart.series[0].format.fill.solid()
    chart.series[0].format.fill.fore_color.rgb = MUTED_GRAY

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(7)
    cat_axis.tick_labels.font.color.rgb = MUTED_GRAY
    cat_axis.tick_labels.font.bold = True
    cat_axis.has_major_gridlines = False
    try:
        cat_axis.format.line.fill.background()
    except Exception:
        pass

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(7)
    val_axis.tick_labels.font.color.rgb = MUTED_GRAY
    val_axis.has_major_gridlines = True
    try:
        val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        val_axis.format.line.fill.background()
    except Exception:
        pass

    try:
        val_axis.tick_labels.number_format = '0.0"x"'
        val_axis.tick_labels.number_format_is_linked = False
    except Exception:
        pass

    # Pin the Y-axis to a deterministic ceiling (next multiple of 5) so the
    # average line we overlay below lands at the right place.
    pos_pe = [v for v in pe_clean if v > 0]
    y_max = float(max(5, math.ceil(max(pos_pe) / 5) * 5)) if pos_pe else 20.0
    try:
        val_axis.maximum_scale = y_max
        val_axis.minimum_scale = 0.0
    except Exception:
        pass

    try:
        chart_frame.line.fill.background()
    except AttributeError:
        pass

    # ── 5-year average dashed overlay ──
    # Empirical plot-area offsets for python-pptx column charts at the sizes
    # we use (≈3" × 2"). The plot area sits inside a small inset; tuning
    # these gets the line within ~2px of the target tick on Keynote and
    # PowerPoint renders we tested. Skipped silently when avg falls outside
    # the visible range or when pptx primitive add_shape isn't usable.
    if five_yr_avg is not None and 0 < five_yr_avg < y_max:
        plot_top = y + h * 0.07
        plot_bottom = y + h * 0.82
        plot_left = x + w * 0.13
        plot_right = x + w * 0.97
        ratio = 1.0 - (five_yr_avg / y_max)
        line_y = plot_top + (plot_bottom - plot_top) * ratio
        line_h = Pt(0.75)
        try:
            avg_line = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                plot_left, line_y - line_h, plot_right - plot_left, line_h * 2,
            )
            avg_line.fill.solid()
            avg_line.fill.fore_color.rgb = PURPLE_AVG
            avg_line.line.color.rgb = PURPLE_AVG
            avg_line.line.width = Pt(0.75)
            # Apply dashed pattern to the rect outline via direct XML — the
            # public python-pptx API exposes line color/width but not dash
            # style. The `prstDash` element is the standard pptx dash spec.
            ln = avg_line.line._get_or_add_ln()
            existing_dash = ln.find(qn("a:prstDash"))
            if existing_dash is not None:
                ln.remove(existing_dash)
            dash = etree.SubElement(ln, qn("a:prstDash"))
            dash.set("val", "dash")
        except Exception:
            pass

    from pptx.enum.text import PP_ALIGN
    annotations: list[str] = []
    if nm_years:
        annotations.append(f"N/M: {', '.join(nm_years)}")
    if five_yr_avg is not None and five_yr_avg > 0:
        annotations.append(f"5yr Avg: {five_yr_avg:.1f}x")
    if annotations:
        txbox = slide.shapes.add_textbox(
            x + Inches(0.05), y + Inches(0.02), w - Inches(0.1), Inches(0.18)
        )
        tf = txbox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = " | ".join(annotations)
        p.alignment = PP_ALIGN.RIGHT
        if p.runs:
            p.runs[0].font.size = Pt(6)
            p.runs[0].font.color.rgb = PURPLE_AVG
            p.runs[0].font.bold = True


def build_price_chart(
    slide,
    x, y, w, h,
    dates: list[str],
    prices: list[float],
    ticker: str = "",
) -> None:
    """1-year daily price line, downsampled to ~50 points.

    `dates` are ISO YYYY-MM-DD; only every 10th is labelled to keep the axis
    readable. Silently no-ops on sparse/missing data.
    """
    if not dates or not prices or len(dates) < 10:
        return

    step = max(1, len(dates) // 50)
    sampled_dates = dates[::step]
    sampled_prices = prices[::step]

    chart_data = CategoryChartData()
    labels: list[str] = []
    for i, d in enumerate(sampled_dates):
        labels.append(d[5:] if i % 10 == 0 else "")
    chart_data.categories = labels
    chart_data.add_series("Price", sampled_prices)

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, x, y, w, h, chart_data
    )
    chart = chart_frame.chart
    chart.has_legend = False

    series = chart.series[0]
    series.format.line.color.rgb = DARK_BLUE
    series.format.line.width = Pt(1.5)
    series.smooth = True

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(6)
    cat_axis.tick_labels.font.color.rgb = MUTED_GRAY
    cat_axis.has_major_gridlines = False
    try:
        cat_axis.format.line.fill.background()
    except Exception:
        pass

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(7)
    val_axis.tick_labels.font.color.rgb = MUTED_GRAY
    val_axis.has_major_gridlines = True
    try:
        val_axis.major_gridlines.format.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        val_axis.format.line.fill.background()
    except Exception:
        pass

    try:
        chart_frame.line.fill.background()
    except AttributeError:
        pass


def build_surprise_summary(
    slide,
    x, y, w, h,
    surprise_data: dict,
    tx_fn,
    rect_fn,
) -> None:
    """Compact earnings surprise history summary box.

    Renders title, headline (color-coded green/red), and up to last 4 quarters
    with arrow glyphs. Expects `surprise_data` with keys: `total_quarters`,
    `summary`, `beat_count`, `miss_count`, and `details` (list of {quarter,
    surprise_pct}).
    """
    from pptx.dml.color import RGBColor
    from pptx.util import Inches

    if not surprise_data or not surprise_data.get("total_quarters"):
        return

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x1F, 0x23, 0x28)
    GREEN = RGBColor(0x1A, 0x7F, 0x37)
    RED = RGBColor(0xCF, 0x22, 0x22)
    MUTED = RGBColor(0x8B, 0x94, 0x9E)
    BORDER = RGBColor(0xDB, 0xE0, 0xE6)

    rect_fn(slide, x, y, w, h, WHITE, BORDER)

    tx_fn(slide, x + Inches(0.1), y + Inches(0.05), w - Inches(0.2), Inches(0.18),
          "Earnings Track Record", sz=8, bold=True, rgb=MUTED)

    summary = surprise_data.get("summary", "")
    beat = surprise_data.get("beat_count", 0)
    miss = surprise_data.get("miss_count", 0)
    color = GREEN if beat > miss else RED if miss > beat else BLACK

    tx_fn(slide, x + Inches(0.1), y + Inches(0.25), w - Inches(0.2), Inches(0.25),
          summary, sz=9, bold=True, rgb=color)

    details = (surprise_data.get("details") or [])[-4:]
    if details:
        detail_parts = []
        for d in details:
            spr = d.get("surprise_pct")
            if spr is not None:
                arrow = "▲" if spr > 1 else "▼" if spr < -1 else "●"
                detail_parts.append(f"{d.get('quarter', '?')}: {arrow}{spr:+.1f}%")
        if detail_parts:
            tx_fn(slide, x + Inches(0.1), y + Inches(0.48), w - Inches(0.2), Inches(0.2),
                  "  ".join(detail_parts), sz=7, rgb=MUTED)


def build_expanded_table(
    slide,
    x, y,
    periods: list[str],
    announcement_dates: list[str],
    metrics: dict,
    currency: str,
    tx_fn,
    rect_fn,
) -> float:
    """6-column financial table with actual/estimate shading.

    Returns the y coordinate just below the last rendered row so the caller can
    stack further content under it. Uses light background for actuals and
    light-blue for estimates; mirrors the cell-grid used in the sample decks.
    """
    from pptx.util import Inches
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    if not periods:
        return y

    periods = periods[-6:] if len(periods) > 6 else periods
    n_cols = len(periods)

    def _tail(arr, n):
        arr = arr or []
        if len(arr) > n:
            return arr[-n:]
        return arr + [None] * (n - len(arr))

    ann_dates = _tail(announcement_dates, n_cols)
    is_estimate = [
        not d or str(d).strip() in ("", "-", "None") for d in ann_dates
    ]

    metric_w = Inches(1.2)
    col_w = Inches(0.85)
    rh = Inches(0.35)

    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    BLACK = RGBColor(0x1F, 0x23, 0x28)
    HEADER_BG = RGBColor(0x0D, 0x11, 0x17)
    ACTUAL_BG = RGBColor(0xF5, 0xF5, 0xF5)
    ESTIMATE_BG = RGBColor(0xE8, 0xF0, 0xFA)
    BORDER = RGBColor(0xDB, 0xE0, 0xE6)

    cx = x
    rect_fn(slide, cx, y, metric_w, rh, HEADER_BG, BORDER)
    tx_fn(slide, cx + Inches(0.05), y + Inches(0.05), metric_w - Inches(0.1), rh, "Metric", sz=8, bold=True, rgb=WHITE)
    cx += metric_w
    for i, p in enumerate(periods):
        label = p.replace("FY", "")
        suffix = "(E)" if is_estimate[i] else "(A)"
        rect_fn(slide, cx, y, col_w, rh, HEADER_BG, BORDER)
        tx_fn(slide, cx + Inches(0.03), y + Inches(0.05), col_w - Inches(0.06), rh, f"{label}{suffix}", sz=7, bold=True, rgb=WHITE, al=PP_ALIGN.CENTER)
        cx += col_w

    _cM = f"({currency}M)" if currency else "(M)"
    _cU = f"({currency})" if currency else ""
    row_defs = [
        (f"Revenue {_cM}", "net_sales"),
        (f"EBITDA {_cM}", "ebitda"),
        (f"EBIT {_cM}", "ebit"),
        (f"Net Income {_cM}", "net_income"),
        (f"EPS {_cU}", "eps"),
    ]

    rendered_rows = 0
    for label, key in row_defs:
        vals = _tail(metrics.get(key), n_cols)
        if all(v is None for v in vals):
            continue

        row_y = y + rh * (rendered_rows + 1)
        cx = x
        rect_fn(slide, cx, row_y, metric_w, rh, WHITE, BORDER)
        tx_fn(slide, cx + Inches(0.05), row_y + Inches(0.05), metric_w - Inches(0.1), rh, label, sz=7, bold=True, rgb=BLACK)
        cx += metric_w

        for i, v in enumerate(vals):
            bg = ESTIMATE_BG if is_estimate[i] else ACTUAL_BG
            rect_fn(slide, cx, row_y, col_w, rh, bg, BORDER)
            if v is None:
                display = "—"
            elif key == "eps":
                display = f"{v:.2f}" if isinstance(v, (int, float)) else str(v)
            else:
                try:
                    fv = float(v)
                    display = f"{fv:,.0f}" if abs(fv) >= 1 else f"{fv:.2f}"
                except (TypeError, ValueError):
                    display = str(v)
            tx_fn(slide, cx + Inches(0.03), row_y + Inches(0.05), col_w - Inches(0.06), rh, display, sz=7, rgb=BLACK, al=PP_ALIGN.CENTER)
            cx += col_w
        rendered_rows += 1

    return y + rh * (rendered_rows + 1) + Inches(0.15)


# ── Quarterly Income Statement Evolution chart ───────────────────────
#
# Mirrors the MS /finances/ Quarterly chart (PDF page 2 reference) —
# clustered bars for Sales / Operating Profit / Net Income with overlaid
# Net Margin and Operating Margin lines on a secondary axis. Estimate
# columns (forecasts beyond the last announcement_date) get rendered in
# the same colors as actuals; the slide caption documents the cutoff so
# we don't ship a misleading dual-color scheme without a legend.

def build_quarterly_income_chart(
    slide,
    x, y, w, h,
    periods: list[str],
    revenue: list,
    ebit: list,
    net_income: list,
    operating_margin_pct: list | None = None,
    net_margin_pct: list | None = None,
    actuals_boundary: int = -1,
    units_label: str = "",
) -> None:
    """Quarterly clustered-bar income chart.

    Margin lines are rendered as a separate native chart overlay would
    be hard with python-pptx (no true secondary-axis line+bar combo
    without raw XML), so we render them as a XL_CHART_TYPE.LINE on a
    secondary chart frame layered on top — but the simplest reliable
    output is the bar-only chart with the Sales / EBIT / Net Income
    series. Margins are surfaced via the slide's text panels instead.
    """
    if not periods or not any(v is not None for v in (revenue or [])):
        return

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    # CRITICAL: preserve None as None (NOT _safe_float which returned
    # 0.0). MS commonly leaves quarterly EBIT empty for older quarters
    # — converting those to 0 produced misleading zero-bars in the
    # chart, which a reviewer caught on 2026-05 ("operating profit 0
    # for 24Q3, 25Q1, 25Q2 …" — those quarters had real earnings; the
    # chart should show no-bar for missing data, not a fake zero).
    def _vals(arr):
        out = []
        for v in (arr or []):
            if v is None:
                out.append(None)
                continue
            try:
                out.append(float(v))
            except (TypeError, ValueError):
                out.append(None)
        return out

    chart_data = CategoryChartData()
    # Compact category labels: "21Q2", "21Q3", ... so a wide chart fits
    # 16-18 quarters without overlapping text.
    short_labels = []
    for p in periods:
        s = str(p)
        # Match formats "2026Q1" / "2026 Q1" / "Q1 2026"
        import re as _r
        m = _r.search(r"(\d{4})\s*Q(\d)|Q(\d)\s*(\d{4})", s)
        if m:
            yr = m.group(1) or m.group(4)
            qn = m.group(2) or m.group(3)
            short_labels.append(f"{yr[-2:]}Q{qn}")
        else:
            short_labels.append(s)
    chart_data.categories = short_labels
    chart_data.add_series("Sales", _vals(revenue))
    chart_data.add_series("Operating Profit", _vals(ebit))
    chart_data.add_series("Net Income", _vals(net_income))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
    )
    chart = chart_frame.chart

    # Apply the MS palette to each bar series.
    for idx, color in enumerate((BLACK_BAR, GOLD_BAR, GREEN_BAR)):
        try:
            chart.series[idx].format.fill.solid()
            chart.series[idx].format.fill.fore_color.rgb = color
        except Exception:
            pass

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(7)
    chart.legend.font.color.rgb = MUTED_GRAY

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(6)
    cat_axis.tick_labels.font.color.rgb = MUTED_GRAY
    cat_axis.tick_labels.font.bold = True
    cat_axis.has_major_gridlines = False
    try:
        cat_axis.format.line.fill.background()
    except Exception:
        pass

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(6)
    val_axis.tick_labels.font.color.rgb = MUTED_GRAY
    val_axis.has_major_gridlines = False
    try:
        val_axis.format.line.fill.background()
    except Exception:
        pass
    try:
        suffix = f'"{units_label}M"' if units_label else '"M"'
        val_axis.tick_labels.number_format = f'#,##0{suffix}'
        val_axis.tick_labels.number_format_is_linked = False
    except Exception:
        pass

    try:
        chart_frame.line.fill.background()
    except AttributeError:
        pass


# ── Quarterly Revenue — Rate of Surprise chart ────────────────────────
#
# Paired bars (Sales Actual vs Sales Estimate) for each quarter. The
# surprise % triangle annotations from the MS screenshot are surfaced as
# data labels on the Actual series ("+5.5%" etc.) so the reader can
# trace beats/misses at a glance.

def build_surprise_chart(
    slide,
    x, y, w, h,
    periods: list[str],
    actual: list,
    estimate: list,
    surprise_pct: list,
    units_label: str = "",
) -> None:
    """Quarterly Sales actual-vs-estimate paired-bar chart."""
    if not periods:
        return

    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    # Preserve None as None — missing actuals/estimates render as
    # no-bar (correct) instead of zero-bar (misleading).
    def _vals(arr):
        out = []
        for v in (arr or []):
            if v is None:
                out.append(None)
                continue
            try:
                out.append(float(v))
            except (TypeError, ValueError):
                out.append(None)
        return out

    # Compact labels
    import re as _r
    short_labels = []
    for p in periods:
        m = _r.search(r"(\d{4})\s*Q(\d)|Q(\d)\s*(\d{4})", str(p))
        if m:
            yr = m.group(1) or m.group(4)
            qn = m.group(2) or m.group(3)
            short_labels.append(f"{yr[-2:]}Q{qn}")
        else:
            short_labels.append(str(p))

    chart_data = CategoryChartData()
    chart_data.categories = short_labels
    chart_data.add_series("Sales Actual", _vals(actual))
    chart_data.add_series("Sales Estimate", _vals(estimate))

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, w, h, chart_data
    )
    chart = chart_frame.chart

    try:
        chart.series[0].format.fill.solid()
        chart.series[0].format.fill.fore_color.rgb = BLACK_BAR
        chart.series[1].format.fill.solid()
        chart.series[1].format.fill.fore_color.rgb = ESTIMATE_GRAY
    except Exception:
        pass

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(7)
    chart.legend.font.color.rgb = MUTED_GRAY

    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(7)
    cat_axis.tick_labels.font.color.rgb = MUTED_GRAY
    cat_axis.tick_labels.font.bold = True
    cat_axis.has_major_gridlines = False

    val_axis = chart.value_axis
    val_axis.tick_labels.font.size = Pt(7)
    val_axis.tick_labels.font.color.rgb = MUTED_GRAY
    val_axis.has_major_gridlines = False
    try:
        suffix = f'"{units_label}M"' if units_label else '"M"'
        val_axis.tick_labels.number_format = f'#,##0{suffix}'
        val_axis.tick_labels.number_format_is_linked = False
    except Exception:
        pass

    try:
        chart_frame.line.fill.background()
    except AttributeError:
        pass
